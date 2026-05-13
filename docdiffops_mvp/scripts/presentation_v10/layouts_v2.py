"""Extended layouts for the v10 enhancement release.

These layouts complement layouts.py and are used by slides_part5.py for:
  * Hero stat slide (massive numbers)
  * Event detail cards (full claim/evidence quotes)
  * Theme cards (per-theme passports)
  * Document spotlights (per-doc passports)
  * Journey timeline (v7-v10 visual)
  * KPI tiles with embedded sparklines
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from .layouts import _add_title_bar, _blank_slide
from .pptx_helpers import (
    add_filled_rect,
    add_footer,
    add_gradient_bg,
    add_image,
    add_picture_centered,
    add_text_box,
    hex_to_rgb,
    set_slide_bg,
)
from .theme import (
    FONT_BODY,
    FONT_HEADER,
    FONT_MONO,
    OCEAN,
    SIZE_BULLET_PT,
    SIZE_FOOTER_PT,
    SIZE_TITLE_PT,
    SLIDE_H_IN,
    SLIDE_W_IN,
    STATUS_HEX,
    STATUS_RU,
    STATUS_TINT_BG,
    TITLE_BAR_H_IN,
)

# Typography constants — available in theme.py Wave-A; safe fallbacks if absent.
SIZE_HERO_STAT_PT: int = getattr(
    __import__("scripts.presentation_v10.theme", fromlist=["SIZE_HERO_STAT_PT"]),
    "SIZE_HERO_STAT_PT",
    120,
)
SIZE_EVENT_QUOTE_PT: int = 14
SIZE_THEME_TITLE_PT: int = 28


# ---------------------------------------------------------------------------
# layout_hero_stat
# ---------------------------------------------------------------------------


def layout_hero_stat(
    prs: Presentation,
    *,
    primary_number: str,
    primary_label: str,
    secondary_numbers: list[str] | None = None,
    secondary_labels: list[str] | None = None,
    caption: str = "",
    section_name: str | None = None,
    page_num: int | None = None,
    toc_slide: object = None,
) -> object:
    """Full-bleed gradient slide with a dominant hero statistic.

    Layout:
      - Gradient background (bg_dark_premium → secondary_premium)
      - Center: primary_number (120pt bronze) + primary_label (24pt white)
      - Optional secondary numbers flanking the primary (left and right)
      - Caption at bottom (14pt italic, muted white)
      - Optional footer when page_num is supplied
    """
    secondary_numbers = secondary_numbers or []
    secondary_labels = secondary_labels or []

    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_dark_premium"])
    add_gradient_bg(slide, OCEAN["bg_dark_premium"], OCEAN["secondary_premium"])

    # ---- Primary number -----------------------------------------------
    # Centered horizontally in the slide, vertically ~40 % from top
    num_w = 8.0
    num_x = (SLIDE_W_IN - num_w) / 2
    num_y = 1.6
    add_text_box(
        slide,
        num_x,
        num_y,
        num_w,
        2.2,
        primary_number,
        font=FONT_HEADER,
        size_pt=SIZE_HERO_STAT_PT,
        bold=True,
        color=OCEAN["bronze"],
        align="center",
    )

    # ---- Primary label ------------------------------------------------
    add_text_box(
        slide,
        num_x,
        num_y + 2.2,
        num_w,
        0.6,
        primary_label,
        font=FONT_BODY,
        size_pt=24,
        bold=False,
        color=OCEAN["ink_inv"],
        align="center",
    )

    # ---- Secondary numbers (up to 3, evenly spaced in the outer thirds) --
    sec_pairs = list(zip(secondary_numbers[:3], secondary_labels[:3]))
    if sec_pairs:
        # Spread them in narrow bands on both sides of center
        positions = [
            (0.3, 2.5),          # far left
            (SLIDE_W_IN - 2.8, 2.5),  # far right
            (0.3, 4.4),          # lower left (3rd item)
        ]
        for i, (val, lbl) in enumerate(sec_pairs):
            px, py = positions[i]
            add_text_box(
                slide, px, py, 2.4, 1.2, val,
                font=FONT_HEADER, size_pt=60, bold=True,
                color=OCEAN["ink_inv"], align="center",
            )
            add_text_box(
                slide, px, py + 1.2, 2.4, 0.4, lbl,
                font=FONT_BODY, size_pt=12, bold=False,
                color="A0C4D8", align="center",
            )

    # ---- Caption -------------------------------------------------------
    if caption:
        add_text_box(
            slide,
            0.5,
            6.45,
            SLIDE_W_IN - 1.0,
            0.45,
            caption,
            font=FONT_BODY,
            size_pt=14,
            bold=False,
            color="A0C4D8",
            align="center",
        )

    # ---- Footer --------------------------------------------------------
    if page_num is not None and section_name is not None:
        add_footer(
            slide,
            section_name=section_name,
            page_num=page_num,
            toc_slide=toc_slide,
        )

    return slide


# ---------------------------------------------------------------------------
# layout_event_card
# ---------------------------------------------------------------------------


def layout_event_card(
    prs: Presentation,
    *,
    event_id: str,
    theme: str,
    status: str,
    claim_text: str,
    claim_source: str,
    evidence_text: str,
    evidence_source: str,
    conclusion: str,
    legal_coordinate: str = "",
    confidence: str = "",
    section_name: str = "События — детали",
    page_num: int,
    toc_slide: object = None,
) -> object:
    """Detailed event card with claim / evidence / conclusion sections.

    Layout (top → bottom):
      Title bar: «E{id} · {theme} · {status}» with status badge
      Claim block  (header strip + italic quote + source)
      Evidence block (header strip + italic quote + source)
      Conclusion + metadata row (legal_coordinate, confidence bar)
    """
    status_hex = STATUS_HEX.get(status, OCEAN["muted"])
    status_ru = STATUS_RU.get(status, status)

    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])

    # ---- Title bar ----------------------------------------------------
    # Omit status from title (status shown as badge); truncate theme to 50 chars
    theme_short = theme[:50] + ("…" if len(theme) > 50 else "")
    title_text = f"{event_id} — {theme_short}"
    _add_title_bar(slide, title_text, size_pt=28.0)

    # Status badge (colored pill on the right side of the title bar)
    badge_w = 1.8
    badge_x = SLIDE_W_IN - badge_w - 0.2
    add_filled_rect(slide, badge_x, 0.1, badge_w, 0.7, status_hex)
    add_text_box(
        slide, badge_x, 0.12, badge_w, 0.66, status_ru,
        font=FONT_BODY, size_pt=11, bold=True,
        color=OCEAN["ink_inv"], align="center",
    )

    # ---- Section blocks (claim / evidence / conclusion) ---------------
    content_x = 0.35
    content_w = SLIDE_W_IN - 0.7
    block_header_h = 0.32
    quote_rule_w = 0.06

    def _block_header(y: float, label: str, bg_hex: str) -> float:
        """Draw a section header strip and return the y after the strip."""
        add_filled_rect(slide, content_x, y, content_w, block_header_h, bg_hex)
        add_text_box(
            slide, content_x + 0.15, y + 0.04, content_w - 0.3, block_header_h - 0.04,
            label, font=FONT_BODY, size_pt=11, bold=True,
            color=OCEAN["ink_inv"], align="left", margin_zero=True,
        )
        return y + block_header_h

    def _quote_box(y: float, body: str, source: str, height: float) -> None:
        """Draw an indented quote block with left rule and source credit."""
        # Left rule
        add_filled_rect(slide, content_x, y, quote_rule_w, height, OCEAN["rule"])
        # Body text
        add_text_box(
            slide, content_x + quote_rule_w + 0.15, y + 0.08,
            content_w - quote_rule_w - 0.2, height - 0.22,
            body, font=FONT_BODY, size_pt=SIZE_EVENT_QUOTE_PT,
            bold=False, color=OCEAN["ink"], align="left",
        )
        # Source credit (muted, small)
        if source:
            add_text_box(
                slide, content_x + quote_rule_w + 0.15, y + height - 0.26,
                content_w - quote_rule_w - 0.2, 0.24,
                source, font=FONT_MONO, size_pt=10,
                bold=False, color=OCEAN["muted"], align="left", margin_zero=True,
            )

    claim_y = 1.0
    claim_body_h = 1.45
    ev_y = claim_y + block_header_h + claim_body_h + 0.1
    ev_body_h = 1.45
    concl_y = ev_y + block_header_h + ev_body_h + 0.12

    y_after = _block_header(claim_y, "УТВЕРЖДЕНИЕ", OCEAN["primary"])
    _quote_box(y_after, claim_text, claim_source, claim_body_h)

    y_after = _block_header(ev_y, "ДОКАЗАТЕЛЬСТВО", OCEAN["secondary"])
    _quote_box(y_after, evidence_text, evidence_source, ev_body_h)

    # ---- Conclusion ---------------------------------------------------
    add_text_box(
        slide, content_x, concl_y, content_w, 0.6,
        conclusion, font=FONT_BODY, size_pt=13,
        bold=False, color=OCEAN["ink"], align="left",
    )

    # Metadata row
    meta_y = concl_y + 0.65
    meta_parts: list[str] = []
    if legal_coordinate:
        meta_parts.append(f"Координата: {legal_coordinate}")
    if confidence:
        meta_parts.append(f"Уверенность: {confidence}")
    if meta_parts:
        add_text_box(
            slide, content_x, meta_y, content_w * 0.75, 0.35,
            "  ·  ".join(meta_parts),
            font=FONT_MONO, size_pt=11, bold=False,
            color=OCEAN["muted"], align="left", margin_zero=True,
        )
        # Mini confidence bar (if confidence is a numeric-ish string)
        try:
            conf_val = float(str(confidence).rstrip("%")) / 100.0
            conf_val = max(0.0, min(1.0, conf_val))
            bar_x = content_x + content_w * 0.75 + 0.1
            bar_w_total = content_w * 0.20
            bar_h = 0.18
            bar_y = meta_y + 0.08
            add_filled_rect(slide, bar_x, bar_y, bar_w_total, bar_h, OCEAN["rule"])
            if conf_val > 0:
                add_filled_rect(slide, bar_x, bar_y,
                                bar_w_total * conf_val, bar_h, status_hex)
        except (ValueError, TypeError):
            pass

    add_footer(slide, section_name=section_name, page_num=page_num, toc_slide=toc_slide)
    return slide


# ---------------------------------------------------------------------------
# layout_theme_card
# ---------------------------------------------------------------------------


def layout_theme_card(
    prs: Presentation,
    *,
    theme_id: str,
    theme_name: str,
    docs_table: list[dict[str, Any]],
    theses: list[str],
    status_breakdown: dict[str, int],
    events_count: int,
    review_count: int,
    section_name: str = "Темы — карточки",
    page_num: int,
    toc_slide: object = None,
) -> object:
    """Per-theme passport card.

    Layout:
      Title bar: «Тема {theme_id}: {theme_name}»
      Left column (5"): docs table (id / code / rank / role)
      Right column: status mini-bars + key theses + event/review footer row
    """
    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])
    _add_title_bar(slide, f"Тема {theme_id}: {theme_name}")

    content_top = 1.05
    left_x = 0.3
    left_w = 4.8
    right_x = left_x + left_w + 0.2
    right_w = SLIDE_W_IN - right_x - 0.3

    # ---- Left column: docs table ------------------------------------
    row_h = 0.32
    col_widths = [0.7, 1.8, 0.5, 1.6]
    col_labels = ["ID", "Код", "Ранг", "Роль"]
    col_keys = ["id", "code", "rank", "role"]

    # Header row
    hdr_y = content_top
    cx = left_x
    for lbl, cw in zip(col_labels, col_widths):
        add_filled_rect(slide, cx, hdr_y, cw - 0.02, row_h, OCEAN["primary"])
        add_text_box(slide, cx + 0.04, hdr_y + 0.04, cw - 0.08, row_h - 0.06,
                     lbl, font=FONT_BODY, size_pt=9, bold=True,
                     color=OCEAN["ink_inv"], align="left", margin_zero=True)
        cx += cw

    # Data rows (up to 13 rows in content area)
    max_rows = 13
    for ri, doc in enumerate(docs_table[:max_rows]):
        ry = content_top + row_h + ri * (row_h + 0.02)
        bg = OCEAN["tile_bg"] if ri % 2 == 0 else OCEAN["bg_light"]
        cx = left_x
        for key, cw in zip(col_keys, col_widths):
            add_filled_rect(slide, cx, ry, cw - 0.02, row_h, bg,
                            line_hex=OCEAN["rule"])
            val = str(doc.get(key, ""))[:18]
            add_text_box(slide, cx + 0.04, ry + 0.04, cw - 0.1, row_h - 0.06,
                         val, font=FONT_MONO, size_pt=8, bold=False,
                         color=OCEAN["ink"], align="left", margin_zero=True)
            cx += cw

    # ---- Right column: status mini-bars -----------------------------
    rcy = content_top
    add_text_box(slide, right_x, rcy, right_w, 0.30, "Статусы",
                 font=FONT_BODY, size_pt=12, bold=True,
                 color=OCEAN["ink"], align="left", margin_zero=True)
    rcy += 0.32

    total_events = max(sum(status_breakdown.values()), 1)
    bar_h = 0.22
    bar_gap = 0.06
    for status, count in sorted(status_breakdown.items(),
                                key=lambda kv: -kv[1]):
        if count == 0:
            continue
        frac = count / total_events
        s_hex = STATUS_HEX.get(status, OCEAN["muted"])
        s_ru = STATUS_RU.get(status, status)
        # Bar background
        add_filled_rect(slide, right_x, rcy, right_w * 0.7, bar_h, OCEAN["rule"])
        # Filled portion
        if frac > 0:
            add_filled_rect(slide, right_x, rcy, right_w * 0.7 * frac, bar_h, s_hex)
        # Label
        add_text_box(slide, right_x + right_w * 0.72, rcy, right_w * 0.28, bar_h,
                     f"{s_ru} {count}", font=FONT_BODY, size_pt=9, bold=False,
                     color=OCEAN["ink"], align="left", margin_zero=True)
        rcy += bar_h + bar_gap

    # ---- Key theses --------------------------------------------------
    rcy += 0.15
    add_text_box(slide, right_x, rcy, right_w, 0.28, "Ключевые тезисы",
                 font=FONT_BODY, size_pt=12, bold=True,
                 color=OCEAN["ink"], align="left", margin_zero=True)
    rcy += 0.30
    for thesis in theses[:3]:
        text = thesis[:100] + ("…" if len(thesis) > 100 else "")
        add_text_box(slide, right_x, rcy, right_w, 0.5,
                     f"▸ {text}", font=FONT_BODY, size_pt=11,
                     bold=False, color=OCEAN["ink"], align="left")
        rcy += 0.54

    # ---- Footer row -------------------------------------------------
    footer_text = f"События: {events_count}  ·  Очередь проверки: {review_count}"
    add_text_box(slide, right_x, 6.55, right_w, 0.32, footer_text,
                 font=FONT_BODY, size_pt=11, bold=False,
                 color=OCEAN["muted"], align="left", margin_zero=True)

    add_footer(slide, section_name=section_name, page_num=page_num, toc_slide=toc_slide)
    return slide


# ---------------------------------------------------------------------------
# layout_doc_passport
# ---------------------------------------------------------------------------


def layout_doc_passport(
    prs: Presentation,
    *,
    doc_id: str,
    doc_code: str,
    doc_title: str,
    doc_metadata: dict[str, Any],
    top_theses: list[str],
    refs_in: list[dict[str, Any]],
    refs_out: list[dict[str, Any]],
    related_events: list[str],
    section_name: str = "Document Spotlights",
    page_num: int,
    toc_slide: object = None,
) -> object:
    """Per-document passport slide.

    Layout:
      Title bar: «{doc_id} {doc_code} — {doc_title}»
      Top metadata strip (6 KPI cells)
      Two-column body:
        Left (5"): top_theses
        Right (7"): refs_in (top) / refs_out (bottom)
      Bottom row: related event IDs
    """
    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])
    doc_title_short = doc_title[:50] + ("…" if len(doc_title) > 50 else "")
    _add_title_bar(slide, f"{doc_id} {doc_code} — {doc_title_short}", size_pt=28.0)

    # ---- Metadata strip (6 KPI cells) --------------------------------
    meta_keys = ["тип", "ранг", "страниц", "URL", "дата", "#событий"]
    meta_vals = [
        str(doc_metadata.get("type", "—")),
        str(doc_metadata.get("rank", "—")),
        str(doc_metadata.get("pages", "—")),
        str(doc_metadata.get("url", "—"))[:25],
        str(doc_metadata.get("date", "—")),
        str(doc_metadata.get("events_count", "—")),
    ]
    strip_y = 0.92
    strip_h = 0.55
    cell_w = SLIDE_W_IN / 6
    for i, (key, val) in enumerate(zip(meta_keys, meta_vals)):
        cx = i * cell_w
        bg = OCEAN["tile_bg"] if i % 2 == 0 else "FFFFFF"
        add_filled_rect(slide, cx, strip_y, cell_w, strip_h, bg,
                        line_hex=OCEAN["rule"])
        add_text_box(slide, cx + 0.05, strip_y + 0.03, cell_w - 0.1, 0.24, key,
                     font=FONT_BODY, size_pt=8, bold=False,
                     color=OCEAN["muted"], align="center", margin_zero=True)
        add_text_box(slide, cx + 0.05, strip_y + 0.27, cell_w - 0.1, 0.25, val,
                     font=FONT_MONO, size_pt=10, bold=True,
                     color=OCEAN["ink"], align="center", margin_zero=True)

    # ---- Two-column body --------------------------------------------
    body_top = strip_y + strip_h + 0.12
    left_x = 0.3
    left_w = 4.8
    right_x = left_x + left_w + 0.25
    right_w = SLIDE_W_IN - right_x - 0.3

    # Left: theses
    add_text_box(slide, left_x, body_top, left_w, 0.28, "Что говорит",
                 font=FONT_BODY, size_pt=12, bold=True,
                 color=OCEAN["ink"], align="left", margin_zero=True)
    ty = body_top + 0.30
    for thesis in top_theses[:3]:
        text = thesis[:90] + ("…" if len(thesis) > 90 else "")
        add_text_box(slide, left_x, ty, left_w, 0.55,
                     f"▸ {text}", font=FONT_BODY, size_pt=11,
                     bold=False, color=OCEAN["ink"], align="left")
        ty += 0.58

    # Right: refs_in
    add_text_box(slide, right_x, body_top, right_w, 0.28, "Кто ссылается",
                 font=FONT_BODY, size_pt=12, bold=True,
                 color=OCEAN["ink"], align="left", margin_zero=True)
    ry = body_top + 0.30
    row_h = 0.30
    for ref in refs_in[:5]:
        ref_str = f"{ref.get('id', '?')}  {ref.get('code', '')[:30]}"
        add_text_box(slide, right_x, ry, right_w, row_h,
                     ref_str, font=FONT_MONO, size_pt=9,
                     bold=False, color=OCEAN["ink"], align="left", margin_zero=True)
        ry += row_h + 0.02

    # Right: refs_out
    ry += 0.1
    add_text_box(slide, right_x, ry, right_w, 0.28, "На кого ссылается",
                 font=FONT_BODY, size_pt=12, bold=True,
                 color=OCEAN["ink"], align="left", margin_zero=True)
    ry += 0.30
    for ref in refs_out[:5]:
        ref_str = f"{ref.get('id', '?')}  {ref.get('code', '')[:30]}"
        add_text_box(slide, right_x, ry, right_w, row_h,
                     ref_str, font=FONT_MONO, size_pt=9,
                     bold=False, color=OCEAN["ink"], align="left", margin_zero=True)
        ry += row_h + 0.02

    # ---- Bottom events row ------------------------------------------
    events_str = "Связанные события: " + "  ·  ".join(related_events[-5:])
    add_text_box(slide, 0.3, 6.50, SLIDE_W_IN - 0.6, 0.36,
                 events_str, font=FONT_MONO, size_pt=10,
                 bold=False, color=OCEAN["muted"], align="left", margin_zero=True)

    add_footer(slide, section_name=section_name, page_num=page_num, toc_slide=toc_slide)
    return slide


# ---------------------------------------------------------------------------
# layout_journey_timeline
# ---------------------------------------------------------------------------


def layout_journey_timeline(
    prs: Presentation,
    *,
    points: list[dict[str, Any]],
    caption: str = "",
    section_name: str,
    page_num: int,
    toc_slide: object = None,
) -> object:
    """Horizontal journey timeline (v7 → v10 or similar).

    If a pre-rendered PNG chart_journey_timeline.png exists in the assets
    directory, it is embedded directly.  Otherwise the timeline is drawn
    natively using circles + text labels.

    points: list of dicts with keys: version / date / docs / pairs / events / note
    """
    from .data_loader import DEFAULT_BUNDLE_DIR  # local import to avoid circular

    assets_dir = DEFAULT_BUNDLE_DIR / "presentation" / "assets"
    png_path = assets_dir / "chart_journey_timeline.png"

    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])
    _add_title_bar(slide, "Путь проекта v7 → v10")

    if png_path.exists():
        # Wave A produced the chart — just embed it
        add_picture_centered(
            slide,
            png_path,
            content_top_in=1.05,
            max_w_in=SLIDE_W_IN - 0.6,
            max_h_in=5.5,
        )
    else:
        # Native fallback: draw the timeline with shapes
        n = len(points)
        if n == 0:
            add_text_box(slide, 0.5, 3.0, SLIDE_W_IN - 1.0, 1.0,
                         "(нет данных для timeline)",
                         font=FONT_BODY, size_pt=14, bold=False,
                         color=OCEAN["muted"], align="center")
        else:
            axis_y = 4.0   # y-position of the horizontal axis line
            axis_x0 = 0.6
            axis_x1 = SLIDE_W_IN - 0.6
            axis_w = axis_x1 - axis_x0

            # Horizontal axis line
            add_filled_rect(slide, axis_x0, axis_y - 0.03, axis_w, 0.06,
                            OCEAN["rule"])

            step = axis_w / max(n - 1, 1)
            circle_r_in = 0.18  # radius in inches (approximate via filled rect)
            for i, pt in enumerate(points):
                cx = axis_x0 + i * step
                # Circle (approximate as small square)
                add_filled_rect(
                    slide,
                    cx - circle_r_in,
                    axis_y - circle_r_in,
                    circle_r_in * 2,
                    circle_r_in * 2,
                    OCEAN["primary"],
                )
                # Version label above axis
                ver = str(pt.get("version", f"v{i}"))
                add_text_box(
                    slide, cx - 0.6, axis_y - 1.05, 1.2, 0.35,
                    ver, font=FONT_HEADER, size_pt=14, bold=True,
                    color=OCEAN["primary"], align="center", margin_zero=True,
                )
                # Date below version
                date_str = str(pt.get("date", ""))
                add_text_box(
                    slide, cx - 0.6, axis_y - 0.65, 1.2, 0.30,
                    date_str, font=FONT_BODY, size_pt=9, bold=False,
                    color=OCEAN["muted"], align="center", margin_zero=True,
                )
                # Metrics below axis
                docs = str(pt.get("docs", ""))
                pairs = str(pt.get("pairs", ""))
                events = str(pt.get("events", ""))
                metrics = f"{docs} docs\n{pairs} pairs\n{events} evts"
                add_text_box(
                    slide, cx - 0.65, axis_y + 0.28, 1.3, 0.80,
                    metrics, font=FONT_MONO, size_pt=8, bold=False,
                    color=OCEAN["ink"], align="center",
                )
                # Note (small, italic-ish)
                note = str(pt.get("note", ""))[:40]
                if note:
                    add_text_box(
                        slide, cx - 0.8, axis_y + 1.15, 1.6, 0.38,
                        note, font=FONT_BODY, size_pt=8, bold=False,
                        color=OCEAN["muted"], align="center", margin_zero=True,
                    )

    if caption:
        add_text_box(
            slide,
            0.5,
            6.45,
            SLIDE_W_IN - 1.0,
            0.42,
            caption,
            font=FONT_BODY,
            size_pt=10,
            bold=False,
            color=OCEAN["muted"],
            align="center",
            margin_zero=True,
        )

    add_footer(slide, section_name=section_name, page_num=page_num, toc_slide=toc_slide)
    return slide


# ---------------------------------------------------------------------------
# layout_kpi_sparkline
# ---------------------------------------------------------------------------


def layout_kpi_sparkline(
    prs: Presentation,
    *,
    title: str,
    tiles: list[dict[str, Any]],
    section_name: str,
    page_num: int,
    toc_slide: object = None,
) -> object:
    """KPI tiles slide — each tile contains a mini-chart in the lower-right corner.

    tiles: list of dicts with keys:
        label           str   — tile label
        value           str   — big number / text displayed prominently
        color           str   — hex color for the value text (no '#')
        microchart_path Path|str|None  — path to a small PNG sparkline image

    Layout mirrors layout_kpi_tiles from layouts.py but embeds the microchart.
    """
    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])
    _add_title_bar(slide, title)

    n = len(tiles)
    if n == 0:
        add_footer(slide, section_name=section_name, page_num=page_num,
                   toc_slide=toc_slide)
        return slide

    content_top = 1.2
    content_h = 5.4
    margin = 0.4

    cols = 2 if n <= 2 else (3 if n <= 6 else 4)
    rows = (n + cols - 1) // cols

    tile_w = (SLIDE_W_IN - 2 * margin - (cols - 1) * 0.2) / cols
    tile_h = (content_h - (rows - 1) * 0.3) / rows

    # Microchart dimensions (embedded in lower-right of tile)
    mc_w_in = 1.5
    mc_h_in = 0.55

    for idx, tile in enumerate(tiles):
        col = idx % cols
        row = idx // cols
        x = margin + col * (tile_w + 0.2)
        y = content_top + row * (tile_h + 0.3)

        label = str(tile.get("label", ""))
        value = str(tile.get("value", ""))
        color_hex = str(tile.get("color", OCEAN["primary"])).lstrip("#")
        mc_path_raw = tile.get("microchart_path")

        # Tile background
        add_filled_rect(slide, x, y, tile_w, tile_h,
                        OCEAN["tile_bg"], line_hex=OCEAN["rule"])

        # Value (big, left-center of tile)
        val_w = tile_w - mc_w_in - 0.15
        add_text_box(
            slide,
            x + 0.12,
            y + 0.20,
            val_w,
            tile_h * 0.55,
            value,
            font=FONT_HEADER,
            size_pt=52,
            bold=True,
            color=color_hex,
            align="left",
        )

        # Label below value
        add_text_box(
            slide,
            x + 0.10,
            y + tile_h * 0.60,
            val_w,
            tile_h * 0.35,
            label,
            font=FONT_BODY,
            size_pt=12,
            bold=False,
            color=OCEAN["ink"],
            align="left",
        )

        # Microchart (lower-right of tile)
        mc_x = x + tile_w - mc_w_in - 0.08
        mc_y = y + tile_h - mc_h_in - 0.08

        if mc_path_raw is not None:
            mc_path = Path(mc_path_raw)
            if mc_path.exists():
                add_image(slide, mc_path, mc_x, mc_y, mc_w_in, mc_h_in)
                continue  # skip placeholder

        # Placeholder rectangle when no microchart PNG is available
        add_filled_rect(slide, mc_x, mc_y, mc_w_in, mc_h_in, OCEAN["rule"])
        add_text_box(
            slide, mc_x, mc_y + 0.12, mc_w_in, mc_h_in - 0.12,
            "▬▬▬",
            font=FONT_BODY, size_pt=10, bold=False,
            color=OCEAN["muted"], align="center", margin_zero=True,
        )

    add_footer(slide, section_name=section_name, page_num=page_num, toc_slide=toc_slide)
    return slide


# ---------------------------------------------------------------------------
# Smoke test
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    from pptx import Presentation as _Prs
    from pptx.util import Inches as _In
    from .data_loader import load_data

    _prs = _Prs()
    _prs.slide_width = _In(13.333)
    _prs.slide_height = _In(7.5)

    _data = load_data()

    # Slide 1: hero stat
    layout_hero_stat(
        _prs,
        primary_number="351",
        primary_label="Пар сравнения",
        secondary_numbers=["27", "312"],
        secondary_labels=["Документов", "Событий"],
        caption="v10 · DocDiffOps — тест hero stat layout",
    )

    # Slide 2: event card
    _e = _data.events_all[0]
    layout_event_card(
        _prs,
        event_id=_e.get("event_id", "E001"),
        theme=_e.get("theme", "—"),
        status=_e.get("status", "match"),
        claim_text=_e.get("claim_left", "") or _e.get("left_text", ""),
        claim_source=_e.get("left_doc", ""),
        evidence_text=_e.get("evidence_right", "") or _e.get("right_text", ""),
        evidence_source=_e.get("right_doc", ""),
        conclusion=_e.get("conclusion", ""),
        legal_coordinate=_e.get("legal_coordinate", ""),
        confidence=_e.get("confidence", ""),
        section_name="Test",
        page_num=2,
        toc_slide=None,
    )

    # Slide 3: journey timeline (native fallback)
    layout_journey_timeline(
        _prs,
        points=[
            {"version": "v7", "date": "2024-Q1", "docs": 18, "pairs": 153, "events": 201, "note": "Baseline"},
            {"version": "v8", "date": "2024-Q3", "docs": 22, "pairs": 231, "events": 278, "note": "Forensic v8"},
            {"version": "v9", "date": "2025-Q1", "docs": 25, "pairs": 300, "events": 295, "note": "Integral"},
            {"version": "v10", "date": "2026-Q2", "docs": 27, "pairs": 351, "events": 312, "note": "Current"},
        ],
        caption="Эволюция корпуса DocDiffOps",
        section_name="Test — Timeline",
        page_num=3,
    )

    _out = Path("/tmp/layouts_v2_test.pptx")
    _prs.save(str(_out))
    print(
        f"Smoke test: {_out} ({_out.stat().st_size // 1024} KB), "
        f"{len(_prs.slides)} slides"
    )
