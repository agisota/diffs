"""Low-level python-pptx builder helpers.

All geometry is expressed in inches; conversion to EMU is done internally.
Every function accepts and returns python-pptx objects so callers can chain
operations.  No business-logic or slide-content here — only primitives.
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .theme import (
    FONT_BODY,
    FONT_HEADER,
    OCEAN,
    SIZE_FOOTER_PT,
    SLIDE_W_IN,
)


# ---------------------------------------------------------------------------
# Color utilities
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert 6-char hex string (without '#') to RGBColor."""
    h = hex_str.lstrip("#")
    r = int(h[0:2], 16)
    g = int(h[2:4], 16)
    b = int(h[4:6], 16)
    return RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Background helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide: object, hex_color: str) -> None:
    """Fill slide background with a solid color."""
    fill = slide.background.fill  # type: ignore[attr-defined]
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def add_gradient_bg(slide: object, hex_start: str, hex_end: str) -> object:
    """Add a horizontal left→right gradient rectangle covering the whole slide.

    Returns the shape so the caller can adjust z-order if needed.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: F401 (just for annotation)

    # Add a full-bleed rectangle (MSO_SHAPE.RECTANGLE = 1)
    shape = slide.shapes.add_shape(  # type: ignore[attr-defined]
        1, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(7.5)
    )
    shape.line.fill.background()  # no outline
    fill = shape.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = hex_to_rgb(hex_start)
    stops[1].position = 1.0
    stops[1].color.rgb = hex_to_rgb(hex_end)
    fill.gradient_angle = 0.0  # 0° = left → right in pptx convention
    return shape


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------

def add_text_box(
    slide: object,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    text: str,
    *,
    font: str = FONT_BODY,
    size_pt: float = 14,
    bold: bool = False,
    color: str = "0F172A",
    align: str = "left",
    margin_zero: bool = False,
) -> object:
    """Add a text box and return the shape.

    Args:
        slide: target slide object
        x_in, y_in, w_in, h_in: position/size in inches
        text: raw string (may contain '\\n' for multi-line)
        font: font family name
        size_pt: font size in points
        bold: bold weight
        color: 6-hex foreground color (no '#')
        align: 'left' | 'center' | 'right'
        margin_zero: zero out internal text box margins
    """
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }

    tb = slide.shapes.add_textbox(  # type: ignore[attr-defined]
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    if margin_zero:
        from pptx.util import Emu
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align_map.get(align, PP_ALIGN.LEFT)
        run = para.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color)

    return tb


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def add_filled_rect(
    slide: object,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    fill_hex: str,
    *,
    line_hex: Optional[str] = None,
) -> object:
    """Add a solid-filled rectangle. Returns the shape."""
    shape = slide.shapes.add_shape(  # type: ignore[attr-defined]
        1,  # MSO_SHAPE.RECTANGLE
        Inches(x_in),
        Inches(y_in),
        Inches(w_in),
        Inches(h_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = hex_to_rgb(line_hex)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()  # no border
    return shape


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------

def add_image(
    slide: object,
    image_path: Path,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
) -> object:
    """Add image at exact position/size. Returns the picture shape."""
    pic = slide.shapes.add_picture(  # type: ignore[attr-defined]
        str(image_path), Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    return pic


def add_picture_centered(
    slide: object,
    image_path: Path,
    *,
    content_top_in: float = 1.4,
    max_w_in: float = 12.33,
    max_h_in: float = 5.5,
) -> object:
    """Add image centered within the content area, preserving aspect ratio."""
    from PIL import Image as PILImage  # type: ignore[import]

    with PILImage.open(str(image_path)) as img:
        img_w, img_h = img.size

    ratio = img_w / img_h
    # Fit within max_w × max_h
    if max_w_in / ratio <= max_h_in:
        w_in = max_w_in
        h_in = max_w_in / ratio
    else:
        h_in = max_h_in
        w_in = max_h_in * ratio

    # Center horizontally, place below content_top_in
    slide_w = SLIDE_W_IN
    x_in = (slide_w - w_in) / 2
    y_in = content_top_in + (max_h_in - h_in) / 2

    return add_image(slide, image_path, x_in, y_in, w_in, h_in)


# ---------------------------------------------------------------------------
# Hyperlink helper
# ---------------------------------------------------------------------------

def set_internal_link(shape: object, target_slide: object) -> None:
    """Set a slide-to-slide hyperlink on *shape*.

    Uses the ``click_action.target_slide`` API available in python-pptx 1.0.2.
    """
    shape.click_action.target_slide = target_slide  # type: ignore[attr-defined]


# ---------------------------------------------------------------------------
# Accessibility helper
# ---------------------------------------------------------------------------

def set_image_alt_text(shape: object, alt_text: str) -> None:
    """Set alt-text on a picture shape via the cNvPr element.

    Required for screen reader / a11y compliance.
    """
    from pptx.oxml.ns import qn
    sp = shape._element  # type: ignore[attr-defined]
    nv_pic_pr = sp.find(qn("p:nvPicPr"))
    if nv_pic_pr is None:
        return
    cNvPr = nv_pic_pr.find(qn("p:cNvPr"))
    if cNvPr is not None:
        cNvPr.set("descr", alt_text)


# ---------------------------------------------------------------------------
# Footer helper
# ---------------------------------------------------------------------------

def add_footer(
    slide: object,
    *,
    section_name: str,
    page_num: int,
    total: int = 153,
    toc_slide: object = None,
) -> None:
    """Add a three-part footer to the slide.

    Left: section_name  |  Center: стр N/total  |  Right: ↑ ToC (linked)
    Footer bar: y = 7.0", h = 0.4"
    """
    footer_y = 7.0
    footer_h = 0.4
    footer_color = OCEAN["muted"]

    # Left — section name
    add_text_box(
        slide,
        0.3,
        footer_y,
        5.0,
        footer_h,
        section_name,
        font=FONT_BODY,
        size_pt=SIZE_FOOTER_PT,
        color=footer_color,
        align="left",
        margin_zero=True,
    )

    # Center — page number
    add_text_box(
        slide,
        5.3,
        footer_y,
        3.0,
        footer_h,
        f"стр {page_num}/{total}",
        font=FONT_BODY,
        size_pt=SIZE_FOOTER_PT,
        color=footer_color,
        align="center",
        margin_zero=True,
    )

    # Right — ToC link
    toc_tb = add_text_box(
        slide,
        10.5,
        footer_y,
        2.5,
        footer_h,
        "↑ ToC",
        font=FONT_BODY,
        size_pt=SIZE_FOOTER_PT,
        color=footer_color,
        align="right",
        margin_zero=True,
    )
    if toc_slide is not None:
        set_internal_link(toc_tb, toc_slide)
