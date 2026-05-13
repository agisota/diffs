"""Layout D — paginated native PPTX table for DocDiffOps v10 presentation.

Public API:
    layout_paginated_table(...) -> tuple[list, int]

Builds N slides, one per page of rows, with zebra stripes, status-colour cells,
word-wrap, and consistent header/footer styling.
"""
from __future__ import annotations

from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .layouts import _add_title_bar, _blank_slide
from .pptx_helpers import add_footer, hex_to_rgb, set_slide_bg
from .theme import (
    FONT_BODY,
    FONT_HEADER,
    OCEAN,
    SIZE_TABLE_BODY_PT,
    SIZE_TABLE_HEADER_PT,
    SLIDE_H_IN,
    SLIDE_W_IN,
    STATUS_TINT_BG,
    STATUS_RU,
    STATUS_TO_MARK,
)


def _set_cell_fill(cell: object, hex_color: str) -> None:
    """Fill a table cell with a solid color."""
    from pptx.oxml.ns import qn
    from lxml import etree

    tc = cell._tc  # type: ignore[attr-defined]
    # Remove any existing fill
    for old in tc.findall(qn("a:solidFill")):
        tc.remove(old)
    # Build solidFill
    solid = etree.SubElement(tc, qn("a:solidFill"))
    srgb = etree.SubElement(solid, qn("a:srgbClr"))
    srgb.set("val", hex_color.lstrip("#").upper())


def _set_cell_text(
    cell: object,
    text: str,
    *,
    font_name: str = FONT_BODY,
    size_pt: float = SIZE_TABLE_BODY_PT,
    bold: bool = False,
    color_hex: str = OCEAN["ink"],
    align: str = "left",
) -> None:
    """Write *text* into a table cell with consistent formatting."""
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    tf = cell.text_frame  # type: ignore[attr-defined]
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align_map.get(align, PP_ALIGN.LEFT)
    # Clear existing runs
    for run in para.runs:
        run.text = ""
    if para.runs:
        run = para.runs[0]
    else:
        run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color_hex)


def layout_paginated_table(
    prs: Presentation,
    *,
    title: str,
    headers: list[str],
    col_widths_in: list[float],
    rows: list[list[str]],
    rows_per_page: int = 15,
    section_name: str,
    page_offset: int,
    total_slides: int = 153,
    toc_slide_idx: int,
    status_col_idx: Optional[int] = None,
    cell_color_map: Optional[dict[str, str]] = None,
    body_font_pt: int = SIZE_TABLE_BODY_PT,
) -> tuple[list, int]:
    """Build N slides of a paginated table.

    Args:
        prs: target Presentation
        title: base title; each slide appends « — стр N из M»
        headers: column header strings
        col_widths_in: column widths in inches (must sum ≤ SLIDE_W_IN - margins)
        rows: list of row data (list of strings per row)
        rows_per_page: max data rows per slide (default 15)
        section_name: footer left text
        page_offset: 1-based slide number for the first page of this table
        total_slides: denominator for footer page numbers
        toc_slide_idx: 0-based index of ToC slide (linked in footer)
        status_col_idx: column index whose value is used for cell fill color
        cell_color_map: raw text → fill hex (overrides STATUS_TINT_BG lookup)

    Returns:
        (slide_objects, n_slides_added)
    """
    # Resolve the ToC slide object for footer wiring
    toc_slide = prs.slides[toc_slide_idx] if toc_slide_idx < len(prs.slides) else None

    # Normalize cell_color_map: merge with STATUS_TINT_BG for status values
    merged_color_map: dict[str, str] = dict(STATUS_TINT_BG)
    if cell_color_map:
        merged_color_map.update(cell_color_map)
    # Also index STATUS_RU → tint (so ru text in cell gets colour too)
    for eng, ru in STATUS_RU.items():
        if eng in STATUS_TINT_BG:
            merged_color_map.setdefault(ru, STATUS_TINT_BG[eng])

    import math
    total_pages = max(1, math.ceil(len(rows) / rows_per_page))

    # Auto-shrink: if the entire dataset fits in one page with rows to spare,
    # use the actual row count so the table doesn't carry empty zebra-stripe
    # padding rows (which look like missing data).
    effective_rpp = rows_per_page
    if total_pages == 1 and 0 < len(rows) < rows_per_page:
        effective_rpp = len(rows)

    slide_objects: list[object] = []

    # Table geometry
    margin_left = 0.35
    margin_right = 0.35
    table_top = 1.0  # just below title bar
    footer_y = 7.15
    table_h = footer_y - table_top - 0.30
    header_row_h = 0.38
    MAX_ROW_HEIGHT_IN = 0.6  # never make rows taller than this (cosmetic)
    data_row_h = min(
        MAX_ROW_HEIGHT_IN,
        (table_h - header_row_h) / max(effective_rpp, 4),
    )

    # Ensure col_widths don't overflow
    total_w = sum(col_widths_in)
    avail_w = SLIDE_W_IN - margin_left - margin_right
    if total_w > avail_w:
        scale = avail_w / total_w
        col_widths_in = [w * scale for w in col_widths_in]
        total_w = avail_w

    for page_idx in range(total_pages):
        page_num = page_idx + 1
        page_rows = rows[page_idx * rows_per_page: (page_idx + 1) * rows_per_page]
        n_data = len(page_rows)
        # Pad to effective_rpp so the table has consistent height. On single-page
        # tables effective_rpp == n_data so no padding is added (no empty stripes).
        while len(page_rows) < effective_rpp:
            page_rows.append([""] * len(headers))

        slide_title = f"{title} — стр {page_num} из {total_pages}"
        slide = _blank_slide(prs)
        set_slide_bg(slide, OCEAN["bg_light"])
        _add_title_bar(slide, slide_title)

        # Create table: 1 header row + effective_rpp data rows
        n_cols = len(headers)
        n_rows_table = 1 + effective_rpp

        actual_table_h = header_row_h + data_row_h * effective_rpp
        table_shape = slide.shapes.add_table(  # type: ignore[attr-defined]
            n_rows_table,
            n_cols,
            Inches(margin_left),
            Inches(table_top),
            Inches(total_w),
            Inches(actual_table_h),
        )
        table = table_shape.table

        # Set column widths
        for ci, cw in enumerate(col_widths_in):
            table.columns[ci].width = Inches(cw)

        # Set row heights
        table.rows[0].height = Inches(header_row_h)
        for ri in range(1, n_rows_table):
            table.rows[ri].height = Inches(data_row_h)

        # Header row
        for ci, hdr in enumerate(headers):
            cell = table.cell(0, ci)
            _set_cell_fill(cell, OCEAN["primary"])
            _set_cell_text(
                cell,
                hdr,
                font_name=FONT_HEADER,
                size_pt=SIZE_TABLE_HEADER_PT,
                bold=True,
                color_hex=OCEAN["ink_inv"],
                align="center",
            )

        # Data rows
        for ri, row_data in enumerate(page_rows):
            is_real = ri < n_data
            even = (ri % 2 == 0)
            row_bg = OCEAN["bg_light"] if even else "F4F8FA"

            for ci, cell_val in enumerate(row_data):
                cell = table.cell(ri + 1, ci)
                display_val = cell_val if is_real else ""
                # Determine fill
                if is_real and status_col_idx is not None and ci == status_col_idx:
                    # Status column: look up by raw English status or Russian label
                    fill_hex = merged_color_map.get(cell_val, row_bg)
                    # Color-blind dual coding: prefix glyph from STATUS_TO_MARK
                    # before status text so status is identifiable without color.
                    glyph = None
                    if cell_val in STATUS_TO_MARK:
                        glyph = STATUS_TO_MARK[cell_val]
                    else:
                        for eng_key, ru_label in STATUS_RU.items():
                            if cell_val == ru_label and eng_key in STATUS_TO_MARK:
                                glyph = STATUS_TO_MARK[eng_key]
                                break
                    if glyph and not display_val.startswith(glyph):
                        display_val = f"{glyph} {display_val}"
                else:
                    fill_hex = row_bg
                _set_cell_fill(cell, fill_hex)
                text_color = OCEAN["ink"] if is_real else OCEAN["bg_light"]
                _set_cell_text(
                    cell,
                    display_val,
                    font_name=FONT_BODY,
                    size_pt=body_font_pt,
                    bold=False,
                    color_hex=text_color,
                    align="left",
                )

        add_footer(
            slide,
            section_name=section_name,
            page_num=page_offset + page_idx,
            total=total_slides,
            toc_slide=toc_slide,
        )
        slide_objects.append(slide)

    return slide_objects, len(slide_objects)
