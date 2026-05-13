"""Helpers for the v10 enhancement release (Wave C).

Called from slides_part1-4.py to insert 32 additional slides:
  insert_hero_stat_slide        →  +1  (slide 7)
  insert_critical_pairs_section →  +2  (divider + table after pair matrix)
  insert_event_detail_cards     → +11  (divider + 10 event cards)
  insert_theme_cards            → +14  (one card per theme T01-T14)
  insert_sankey_treemap_slides  →  +2  (chart slides in themes section, net 0 after removing 2)
  insert_journey_timeline_slide →  +1  (in Trend & QA section)
  insert_doc_spotlights         →  +3  (D18, D24, D27 at end)
"""
from __future__ import annotations

from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

from .data_loader import V10Data
from .layouts import layout_divider, layout_full_chart, layout_bullets_visual
from .layouts_v2 import (
    layout_hero_stat,
    layout_event_card,
    layout_theme_card,
    layout_doc_passport,
    layout_journey_timeline,
)
from .layout_table import layout_paginated_table
from .theme import OCEAN, STATUS_HEX, STATUS_RU, STATUS_TINT_BG, V8_STATUSES

REPO_ROOT = Path(__file__).resolve().parents[3]
ASSETS_DIR = REPO_ROOT / "migration_v10_out" / "presentation" / "assets"


# ---------------------------------------------------------------------------
# Hero stat slide (inserted as slide 7)
# ---------------------------------------------------------------------------


def insert_hero_stat_slide(
    prs: Presentation,
    data: V10Data,
    *,
    page_num: int,
    toc_slide: object,
) -> None:
    """Insert a massive hero-stat slide.  «1 / 1 / 351» — совпадение / противоречие / пары."""
    pairs_status = data.pairs_by_status()
    layout_hero_stat(
        prs,
        primary_number=str(pairs_status.get("contradiction", 1)),
        primary_label="противоречие",
        secondary_numbers=[
            str(pairs_status.get("match", 1)),
            str(sum(pairs_status.values())),
        ],
        secondary_labels=["совпадение", "всего пар"],
        caption=(
            "итог попарного сравнения 27 нормативных и аналитических"
            " документов миграционной политики РФ"
        ),
        section_name="Executive Summary",
        page_num=page_num,
        toc_slide=toc_slide,
    )


# ---------------------------------------------------------------------------
# Critical pairs section  (+2 slides: divider + 1-page table)
# ---------------------------------------------------------------------------


def insert_critical_pairs_section(
    prs: Presentation,
    data: V10Data,
    *,
    page_offset: int,
    toc_slide: object,
) -> int:
    """Add divider + critical-pairs table.  Returns 0-based index of the divider slide."""
    layout_divider(
        prs,
        section_no="03·1",
        section_title="Критические пары: contradiction + partial_overlap",
        toc_slide=toc_slide,
    )
    div_idx = len(prs.slides) - 1

    critical = [
        p for p in data.pairs
        if p.get("v8_status") in ("contradiction", "partial_overlap")
    ]

    rows: list[list[str]] = []
    for p in critical:
        left_short = _trunc(data.doc_short(p.get("left", "")), 20)
        right_short = _trunc(data.doc_short(p.get("right", "")), 20)
        status_ru = STATUS_RU.get(p.get("v8_status", ""), p.get("v8_status", ""))
        rows.append([
            p.get("id", ""),
            left_short,
            right_short,
            _trunc(p.get("topics", ""), 38),
            status_ru,
            str(p.get("events_count", "0")),
            p.get("rank_pair", ""),
        ])

    layout_paginated_table(
        prs,
        title="Критические пары: contradiction + partial_overlap",
        headers=["ID", "Левый", "Правый", "Темы", "Статус", "Соб.", "Ранги"],
        col_widths_in=[1.0, 1.8, 1.8, 3.0, 1.8, 0.7, 1.2],
        rows=rows,
        rows_per_page=27,
        section_name="Критические пары",
        page_offset=page_offset + 1,
        total_slides=153,
        toc_slide_idx=2,
        status_col_idx=4,
        cell_color_map=STATUS_TINT_BG,
    )
    return div_idx


# ---------------------------------------------------------------------------
# Event detail cards (+11 slides: divider + 10 cards)
# ---------------------------------------------------------------------------


def _select_top_events(data: V10Data) -> list[dict[str, Any]]:
    """1 contradiction + 6 partial_overlap + 3 outdated → 10 total."""
    by_status: dict[str, list[dict[str, Any]]] = {s: [] for s in V8_STATUSES}
    for ev in data.events_all:
        s = ev.get("status", "")
        if s in by_status:
            by_status[s].append(ev)

    # sort each bucket by confidence desc
    for s in by_status:
        by_status[s].sort(
            key=lambda e: float(e.get("confidence", "0") or "0"),
            reverse=True,
        )

    selected: list[dict[str, Any]] = []
    selected.extend(by_status["contradiction"][:1])
    selected.extend(by_status["partial_overlap"][:6])
    selected.extend(by_status["outdated"][:3])
    return selected[:10]


def insert_event_detail_cards(
    prs: Presentation,
    data: V10Data,
    *,
    page_offset: int,
    toc_slide: object,
) -> tuple[int, dict[str, int]]:
    """Add divider + 10 event cards.

    Returns (div_idx, {event_id: slide_0based_idx}).
    """
    layout_divider(
        prs,
        section_no="04·1",
        section_title="События — детальные карточки (1 противоречие · 6 частичных · 3 устаревших)",
        toc_slide=toc_slide,
    )
    div_idx = len(prs.slides) - 1

    events = _select_top_events(data)
    indices: dict[str, int] = {}

    for i, ev in enumerate(events):
        layout_event_card(
            prs,
            event_id=ev.get("event_id", ""),
            theme=ev.get("theme", ""),
            status=ev.get("status", ""),
            claim_text=_trunc(ev.get("claim_left", ""), 300),
            claim_source=f"{ev.get('left_id', '')} {ev.get('left_doc', '')}",
            evidence_text=_trunc(ev.get("evidence_right", ""), 300),
            evidence_source=f"{ev.get('right_id', '')} {ev.get('right_doc', '')}",
            conclusion=_trunc(ev.get("conclusion", ""), 200),
            legal_coordinate=ev.get("legal_coordinate", ""),
            confidence=ev.get("confidence", ""),
            section_name="События — детали",
            page_num=page_offset + 1 + i,
            toc_slide=toc_slide,
        )
        indices[ev.get("event_id", "")] = len(prs.slides) - 1

    return div_idx, indices


# ---------------------------------------------------------------------------
# Theme cards (+14 slides)
# ---------------------------------------------------------------------------


def _select_core_themes(data: V10Data, n: int = 14) -> list[dict[str, Any]]:
    """Build theme meta sorted by event count desc, return top n."""
    theme_meta: dict[str, dict[str, Any]] = {}
    for row in data.theme_doc:
        tid = row.get("theme_id", "")
        if not tid:
            continue
        if tid not in theme_meta:
            theme_meta[tid] = {
                "id": tid,
                "name": row.get("theme", ""),
                "doc_count": 0,
                "event_count": 0,
            }
        theme_meta[tid]["doc_count"] += 1

    for ev in data.events_all:
        tid = ev.get("theme_id", "")
        if tid in theme_meta:
            theme_meta[tid]["event_count"] += 1

    candidates = list(theme_meta.values())
    candidates.sort(key=lambda t: -t["event_count"])
    return candidates[:n]


def insert_theme_cards(
    prs: Presentation,
    data: V10Data,
    *,
    page_offset: int,
    toc_slide: object,
) -> dict[str, int]:
    """Add 14 theme-card slides.  Returns {theme_id: slide_0based_idx}."""
    themes = _select_core_themes(data, n=14)
    indices: dict[str, int] = {}

    for i, theme in enumerate(themes):
        tid = theme["id"]

        # Docs in this theme
        docs_in_theme: list[dict[str, Any]] = []
        for row in data.theme_doc:
            if row.get("theme_id") != tid:
                continue
            doc_id = row.get("doc_id", "")
            d = data.doc_by_id(doc_id)
            if d:
                docs_in_theme.append({
                    "id": doc_id,
                    "code": _trunc(d.get("code", ""), 18),
                    "rank": d.get("rank", ""),
                    "role": _trunc(row.get("role", ""), 16),
                })

        # Status breakdown for events in this theme
        status_breakdown: Counter[str] = Counter()
        for ev in data.events_all:
            if ev.get("theme_id") == tid:
                status_breakdown[ev.get("status", "")] += 1

        # Top theses: match by theme name (theses have no theme_id col)
        theses_for_theme = [
            t for t in data.theses if t.get("theme", "") == theme["name"]
        ][:3]
        thesis_texts = [
            _trunc(t.get("thesis", t.get("claim_text", "")), 100)
            for t in theses_for_theme
        ]

        events_count = sum(status_breakdown.values())
        review_count = sum(
            1 for r in data.review_queue
            if r.get("theme", "") == theme["name"]
        )

        layout_theme_card(
            prs,
            theme_id=tid,
            theme_name=theme["name"],
            docs_table=docs_in_theme[:8],
            theses=thesis_texts,
            status_breakdown=dict(status_breakdown),
            events_count=events_count,
            review_count=review_count,
            section_name="Темы — карточки",
            page_num=page_offset + i,
            toc_slide=toc_slide,
        )
        indices[tid] = len(prs.slides) - 1

    return indices


# ---------------------------------------------------------------------------
# Sankey + Treemap slides (2 chart slides for the themes section)
# ---------------------------------------------------------------------------


def insert_sankey_treemap_slides(
    prs: Presentation,
    *,
    page_num_start: int,
    toc_slide: object,
) -> None:
    """Add 2 chart slides: Sankey rank-flow + Treemap themes."""
    layout_full_chart(
        prs,
        title="Поток статусов через ранги документов",
        image_path=ASSETS_DIR / "chart_sankey_rank_flow.png",
        caption=(
            "Sankey: rank-pair → v8 status."
            " Толщина потока = число пар."
        ),
        section_name="Темы и корреляция",
        page_num=page_num_start,
        toc_slide=toc_slide,
    )
    layout_full_chart(
        prs,
        title="Темы по плотности событий",
        image_path=ASSETS_DIR / "chart_treemap_themes.png",
        caption=(
            "Treemap: площадь = число событий,"
            " цвет = доминирующий статус."
        ),
        section_name="Темы и корреляция",
        page_num=page_num_start + 1,
        toc_slide=toc_slide,
    )


# ---------------------------------------------------------------------------
# Journey timeline (+1 slide)
# ---------------------------------------------------------------------------


def insert_journey_timeline_slide(
    prs: Presentation,
    data: V10Data,
    *,
    page_num: int,
    toc_slide: object,
) -> None:
    """Add one journey-timeline slide (v7→v10 annotated milestones)."""
    timeline = data.trend.get("timeline", [])
    points: list[dict[str, Any]] = []
    for t in timeline:
        points.append({
            "version": t.get("version", ""),
            "date": str(t.get("date", ""))[:10],
            "docs": t.get("docs", 0),
            "pairs": t.get("pairs", 0),
            "events": t.get("events", 0),
            "note": t.get("note", ""),
        })

    layout_journey_timeline(
        prs,
        points=points,
        caption="v9 — расширение корпуса, не регресс качества",
        section_name="Тренд и QA",
        page_num=page_num,
        toc_slide=toc_slide,
    )


# ---------------------------------------------------------------------------
# Document spotlights (+3 slides: D18, D24, D27)
# ---------------------------------------------------------------------------


def _select_spotlight_docs(data: V10Data) -> list[dict[str, str]]:
    target_ids = ("D18", "D24", "D27")
    ordered = {doc_id: None for doc_id in target_ids}
    for d in data.documents:
        if d.get("id") in target_ids:
            ordered[d["id"]] = d  # type: ignore[assignment]
    return [v for v in ordered.values() if v is not None]


def insert_doc_spotlights(
    prs: Presentation,
    data: V10Data,
    *,
    page_offset: int,
    toc_slide: object,
) -> dict[str, int]:
    """Add 3 document-passport slides for D18, D24, D27.

    Returns {doc_id: slide_0based_idx}.
    """
    docs = _select_spotlight_docs(data)
    indices: dict[str, int] = {}

    for i, d in enumerate(docs):
        doc_id = d.get("id", "")

        # Top 3 theses from theses CSV (matching by source_doc code)
        doc_code = d.get("code", "")
        theses_for_doc = [
            t for t in data.theses
            if t.get("source_doc", "") == doc_code
        ][:3]
        top_theses = [
            _trunc(t.get("thesis", t.get("claim_text", "")), 120)
            for t in theses_for_doc
        ]

        # refs in / out from dependency_graph
        refs_out_raw = [
            e for e in data.dependency_graph
            if e.get("from_doc_id") == doc_id
        ][:5]
        refs_in_raw = [
            e for e in data.dependency_graph
            if e.get("to_doc_id") == doc_id
        ][:5]

        # related events
        related_events = [
            ev.get("event_id", "")
            for ev in data.events_all
            if ev.get("left_id") == doc_id or ev.get("right_id") == doc_id
        ][:5]

        layout_doc_passport(
            prs,
            doc_id=doc_id,
            doc_code=d.get("code", ""),
            doc_title=_trunc(d.get("title", ""), 60),
            doc_metadata={
                "type": _trunc(d.get("type", ""), 30),
                "rank": d.get("rank", ""),
                "url": _trunc(d.get("url", ""), 50),
            },
            top_theses=top_theses,
            refs_in=[
                {"id": e.get("from_doc_id", ""), "code": e.get("from_doc_short", "")}
                for e in refs_in_raw
            ],
            refs_out=[
                {"id": e.get("to_doc_id", ""), "code": e.get("to_doc_short", "")}
                for e in refs_out_raw
            ],
            related_events=related_events,
            section_name="Document Spotlights",
            page_num=page_offset + i,
            toc_slide=toc_slide,
        )
        indices[doc_id] = len(prs.slides) - 1

    return indices


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _trunc(text: str, n: int) -> str:
    if len(text) <= n:
        return text
    return text[: n - 1] + "…"
