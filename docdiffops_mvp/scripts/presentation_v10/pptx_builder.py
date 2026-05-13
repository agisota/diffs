"""Main PPTX builder entry-point for DocDiffOps v10 presentation.

Usage:
    cd docdiffops_mvp
    python -m scripts.presentation_v10.pptx_builder

Output:
    migration_v10_out/presentation/DocDiffOps_v10_presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from .data_loader import load_data, V10Data
from .slides_part1 import build_part1
from .slides_part2 import build_part2
from .slides_part3 import build_part3
from .slides_part4 import build_part4
from .cross_links import rewire_all as cross_link_all
from .pptx_helpers import set_internal_link
from .speaker_notes import attach_notes_to_slides
from .theme import SLIDE_H_IN, SLIDE_W_IN

REPO_ROOT = Path(__file__).resolve().parents[3]
DEFAULT_OUT = REPO_ROOT / "migration_v10_out" / "presentation" / "DocDiffOps_v10_presentation.pptx"

# Mapping: ToC tile index → refs key for the section divider
# Tiles in slide 3 are ordered: 0=Executive, 1=Корпус, 2=Матрица пар,
# 3=События, 4=Темы, 5=Очередь проверки, 6=Тренд и QA, 7=Действия, 8=↑Обложка
_TOC_TILE_REFS = [
    "executive",        # tile 0 → section 01
    "corpus_div",       # tile 1 → section 02
    "pair_matrix_div",  # tile 2 → section 03
    "events_div",       # tile 3 → section 04
    "themes_div",       # tile 4 → section 05
    "review_div",       # tile 5 → section 06
    "trend_qa_div",     # tile 6 → section 07
    "actions_div",      # tile 7 → section 08
    "cover",            # tile 8 → cover slide
]


def _rewire_toc(prs: Presentation, refs: dict) -> None:
    """Update ToC tile hyperlinks to point to actual section divider slides.

    parts 2-4 populate refs with the real 0-based indices after their
    divider slides are created.  We iterate the tile shapes stored on the
    ToC slide and set correct target_slide for each.
    """
    toc_slide = prs.slides[refs["toc"]]
    tile_shapes = getattr(toc_slide, "_toc_tile_shapes", [])
    for tile_idx, ref_key in enumerate(_TOC_TILE_REFS):
        if tile_idx >= len(tile_shapes):
            break
        slide_idx = refs.get(ref_key)
        if slide_idx is None:
            continue
        if slide_idx < len(prs.slides):
            set_internal_link(tile_shapes[tile_idx], prs.slides[slide_idx])


def build_pptx(out_path: Path, *, data: V10Data | None = None) -> Path:
    """Build the presentation and save to *out_path*. Returns the path."""
    if data is None:
        data = load_data()

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    refs: dict[str, int] = {}
    build_part1(prs, data, refs=refs)
    build_part2(prs, data, refs=refs)
    build_part3(prs, data, refs=refs)
    build_part4(prs, data, refs=refs)

    # Rewire ToC links now that all dividers exist
    _rewire_toc(prs, refs)

    # Cross-links event↔pair↔doc (after all slides built and ToC wired)
    link_counts = cross_link_all(prs)
    print(f"  Cross-links: {link_counts['total_links']} total "
          f"(events: {link_counts['events_table_links']}, "
          f"review: {link_counts['review_queue_links']}, "
          f"FA: {link_counts['fa_doc_links']})")

    # Speaker notes on all slides (after all slides and cross-links are set)
    n_notes = attach_notes_to_slides(prs, data)
    print(f"  Speaker notes: {n_notes} slides annotated")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


if __name__ == "__main__":
    out = build_pptx(DEFAULT_OUT)
    size_kb = out.stat().st_size // 1024
    slide_count = len(Presentation(str(out)).slides)
    print(f"Saved: {out} ({size_kb} KB), slides={slide_count}")
