"""Cross-link registry and patcher for the v10 presentation.

After all slides are built, this module walks tables in events / review_queue /
FA cards and patches each cell containing an event_id or doc_id with a
slide-internal hyperlink to the corresponding detail card / spotlight.

Called from pptx_builder.build_pptx() AFTER all parts are built and ToC is wired.

Note on python-pptx API limitations
------------------------------------
``_Run.click_action`` is not available on table-cell runs or text-box runs in
python-pptx 1.0.2.  We therefore inject the OOXML ``<a:hlinkClick>`` element
directly into the run's ``<a:rPr>`` node and create the required slide
relationship via ``slide.part.relate_to()``.  This is the standard PPTX
mechanism for in-presentation slide jumps and is well-supported by PowerPoint,
LibreOffice Impress, and Apple Keynote.
"""
from __future__ import annotations

import re
from typing import Iterable

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from .pptx_helpers import hex_to_rgb
from .theme import OCEAN

# ---------------------------------------------------------------------------
# Patterns
# ---------------------------------------------------------------------------

# Event IDs: E001 .. E999
_EVENT_ID_RE = re.compile(r"\b(E\d{3})\b")
# Doc IDs: D01 .. D99
_DOC_ID_RE = re.compile(r"\b(D\d{2})\b")

# Relationship type for slide→slide hyperlink
_SLIDE_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
)


# ---------------------------------------------------------------------------
# Low-level XML helper
# ---------------------------------------------------------------------------

def _inject_slide_link(run: object, source_slide: object, target_slide: object) -> None:
    """Inject an <a:hlinkClick> slide-jump hyperlink into *run*.

    Works for both table-cell runs and textbox runs by manipulating OOXML
    directly, since python-pptx 1.0.2 does not expose ``click_action`` on
    ``_Run`` objects.

    Also sets font colour to OCEAN["primary"] and underline=True so the link
    is visually discoverable.
    """
    # Create (or reuse) slide→target relationship
    source_part = source_slide.part  # type: ignore[attr-defined]
    target_part = target_slide.part  # type: ignore[attr-defined]
    rId = source_part.relate_to(target_part, _SLIDE_REL_TYPE)

    # Get underlying lxml element for the run
    r_elem = run._r  # type: ignore[attr-defined]

    # Ensure <a:rPr> exists as first child
    rPr = r_elem.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(r_elem, qn("a:rPr"))
        r_elem.insert(0, rPr)

    # Remove any pre-existing hlinkClick to avoid duplicates
    for old in rPr.findall(qn("a:hlinkClick")):
        rPr.remove(old)

    # Inject hlinkClick with slide-jump action
    hlinkClick = etree.SubElement(rPr, qn("a:hlinkClick"))
    hlinkClick.set(qn("r:id"), rId)
    hlinkClick.set("action", "ppaction://hlinksldjump")

    # Visual hint: ocean-blue + underline
    try:
        rgb = hex_to_rgb(OCEAN["primary"])
        # Apply colour: solidFill inside rPr
        # Remove existing solidFill first
        for old_fill in rPr.findall(qn("a:solidFill")):
            rPr.remove(old_fill)
        solid_fill = etree.SubElement(rPr, qn("a:solidFill"))
        srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
        srgb.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
        # Underline attribute on rPr
        rPr.set("u", "sng")
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Collect refs from built slides
# ---------------------------------------------------------------------------

def collect_event_card_refs(
    prs: object,
    *,
    event_card_start_idx: int = 76,
    event_card_end_idx: int = 86,
) -> dict[str, int]:
    """Walk slides in [event_card_start_idx, event_card_end_idx] and return
    {event_id: slide_0based_idx} by reading each slide's title-like shapes.

    The first shape whose text matches E\\d{3} determines the event_id for
    that slide.
    """
    out: dict[str, int] = {}
    slides = prs.slides  # type: ignore[attr-defined]
    for idx in range(event_card_start_idx, min(event_card_end_idx + 1, len(slides))):
        slide = slides[idx]
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            text = shape.text_frame.text or ""
            m = _EVENT_ID_RE.search(text)
            if m:
                out[m.group(1)] = idx
                break
    return out


def collect_doc_spotlight_refs(
    prs: object,
    *,
    spotlight_start_idx: int = 149,
    spotlight_end_idx: int = 151,
) -> dict[str, int]:
    """Walk slides in [spotlight_start_idx, spotlight_end_idx] and return
    {doc_id: slide_0based_idx} by reading each slide's title-like shapes.
    """
    out: dict[str, int] = {}
    slides = prs.slides  # type: ignore[attr-defined]
    for idx in range(spotlight_start_idx, min(spotlight_end_idx + 1, len(slides))):
        slide = slides[idx]
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            text = shape.text_frame.text or ""
            m = _DOC_ID_RE.search(text)
            if m:
                out[m.group(1)] = idx
                break
    return out


# ---------------------------------------------------------------------------
# Patch events table and review queue
# ---------------------------------------------------------------------------

def patch_table_cells_with_event_links(
    prs: object,
    *,
    slide_indices: Iterable[int],
    event_id_col: int,
    event_card_refs: dict[str, int],
) -> int:
    """In each table on each specified slide, find cells in *event_id_col*
    that contain an event_id and inject a slide-jump hyperlink.

    Returns the count of links successfully injected.
    """
    added = 0
    slides = prs.slides  # type: ignore[attr-defined]
    for slide_idx in slide_indices:
        if slide_idx >= len(slides):
            continue
        slide = slides[slide_idx]
        for shape in slide.shapes:
            if not hasattr(shape, "table"):
                continue
            table = shape.table
            num_cols = len(table.columns)
            if event_id_col >= num_cols:
                continue
            # Row 0 is header; data starts at row 1
            for row_idx in range(1, len(table.rows)):
                cell = table.cell(row_idx, event_id_col)
                if not hasattr(cell, "text_frame"):
                    continue
                cell_text = cell.text_frame.text.strip()
                m = _EVENT_ID_RE.search(cell_text)
                if not m:
                    continue
                event_id = m.group(1)
                target_idx = event_card_refs.get(event_id)
                if target_idx is None:
                    continue
                target_slide = slides[target_idx]
                # Inject link on first run of first paragraph
                tf = cell.text_frame
                if not tf.paragraphs:
                    continue
                para = tf.paragraphs[0]
                if not para.runs:
                    # Add a run if the cell text is set without a proper run
                    run = para.add_run()
                    run.text = cell_text
                else:
                    run = para.runs[0]
                try:
                    _inject_slide_link(run, slide, target_slide)
                    added += 1
                except Exception:
                    pass
    return added


# ---------------------------------------------------------------------------
# Patch FA slides with doc links
# ---------------------------------------------------------------------------

def patch_fa_slides_with_doc_links(
    prs: object,
    *,
    fa_slide_indices: Iterable[int],
    doc_spotlight_refs: dict[str, int],
) -> int:
    """In FA slide text boxes, find runs containing D## doc IDs and inject
    slide-jump hyperlinks to the corresponding doc spotlight slide.

    Returns the count of links successfully injected.
    """
    added = 0
    slides = prs.slides  # type: ignore[attr-defined]
    for slide_idx in fa_slide_indices:
        if slide_idx >= len(slides):
            continue
        slide = slides[slide_idx]
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    text = run.text or ""
                    m = _DOC_ID_RE.search(text)
                    if not m:
                        continue
                    doc_id = m.group(1)
                    target_idx = doc_spotlight_refs.get(doc_id)
                    if target_idx is None:
                        continue
                    target_slide = slides[target_idx]
                    try:
                        _inject_slide_link(run, slide, target_slide)
                        added += 1
                    except Exception:
                        pass
    return added


# ---------------------------------------------------------------------------
# Main entry-point
# ---------------------------------------------------------------------------

def rewire_all(prs: object) -> dict[str, int]:
    """Run cross-link patching after all slides are built.

    Slide index ranges (0-based) follow the v10 layout:
      - Events table:    slides 53-74  (1-based 54-75), col 0 = event_id
      - Event cards:     slides 76-86  (1-based 77-87), title contains E###
      - Review queue:    slides 109-122 (1-based 110-123), col 2 = event_id
      - FA cards:        slides 137-146 (1-based 138-147), bullets may ref D##
      - Doc spotlights:  slides 149-151 (1-based 150-152), title contains D##

    Returns a counts dict for logging.
    """
    # Collect mappings from already-built slides
    event_refs = collect_event_card_refs(prs, event_card_start_idx=76, event_card_end_idx=86)
    doc_refs = collect_doc_spotlight_refs(prs, spotlight_start_idx=149, spotlight_end_idx=151)

    # Patch events table (slides 53-74, col 0 = event_id)
    events_links = patch_table_cells_with_event_links(
        prs,
        slide_indices=range(53, 75),
        event_id_col=0,
        event_card_refs=event_refs,
    )

    # Patch review queue (slides 109-122, col 2 = event_id)
    review_links = patch_table_cells_with_event_links(
        prs,
        slide_indices=range(109, 123),
        event_id_col=2,
        event_card_refs=event_refs,
    )

    # Patch FA card slides (137-146) for doc references
    fa_links = patch_fa_slides_with_doc_links(
        prs,
        fa_slide_indices=range(137, 147),
        doc_spotlight_refs=doc_refs,
    )

    return {
        "event_card_refs": len(event_refs),
        "doc_spotlight_refs": len(doc_refs),
        "events_table_links": events_links,
        "review_queue_links": review_links,
        "fa_doc_links": fa_links,
        "total_links": events_links + review_links + fa_links,
    }
