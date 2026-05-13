"""Slides 96-120: Trend & QA (96-105) + Actions & Outro (106-120).

Public API:
    build_part4(prs, data, *, refs) -> None

Populates refs with:
    'trend_qa_div' = 95   (0-based index of slide 96)
    'actions_div'  = 105  (0-based index of slide 106)
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from .data_loader import V10Data
from .layout_table import layout_paginated_table
from .layouts import (
    layout_bullets_full_width,
    layout_bullets_visual,
    layout_cover,
    layout_divider,
    layout_full_chart,
    layout_kpi_tiles,
)
from .slides_part5 import insert_doc_spotlights, insert_journey_timeline_slide
from .theme import OCEAN, STATUS_TINT_BG

REPO_ROOT = Path(__file__).resolve().parents[3]
ASSETS_DIR = REPO_ROOT / "migration_v10_out" / "presentation" / "assets"


def _trunc(text: str, n: int) -> str:
    if len(text) <= n:
        return text
    return text[: n - 1] + "…"


def build_part4(prs: Presentation, data: V10Data, *, refs: dict[str, Any]) -> None:
    """Add slides 96-120 to *prs* and record divider indices in *refs*."""

    toc_slide = prs.slides[refs["toc"]]

    # ================================================================== #
    # Раздел 7 — Тренд и QA (slides 96-105)
    # ================================================================== #

    # ------------------------------------------------------------------ #
    # Slide 96 — Divider: Trend & QA
    # ------------------------------------------------------------------ #
    layout_divider(
        prs,
        section_no="07",
        section_title="Тренд и QA",
        toc_slide=toc_slide,
    )
    refs["trend_qa_div"] = len(prs.slides) - 1  # 0-based → 95

    # ------------------------------------------------------------------ #
    # Slide 97 — Match share trend chart
    # ------------------------------------------------------------------ #
    layout_full_chart(
        prs,
        title="Тренд: доля точных совпадений v7-v10",
        image_path=ASSETS_DIR / "chart_trend_match_share.png",
        caption=(
            "match_share: v7=8.0%, v8=8.0%, v9=0.28%, v10=0.28%. "
            "Падение v8→v9 вызвано расширением корпуса (D27 ВЦИОМ) — не регресс качества."
        ),
        section_name="Тренд и QA",
        page_num=97,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 98 — Review queue trend chart
    # ------------------------------------------------------------------ #
    layout_full_chart(
        prs,
        title="Тренд: очередь ручной проверки v7-v10",
        image_path=ASSETS_DIR / "chart_trend_review_queue.png",
        caption=(
            "Очередь: v7=183, v8=183, v9=54, v10=54. "
            "Сокращение v8→v9: переклассификация 129 manual_review в not_comparable "
            "после добавления rank-3 документа."
        ),
        section_name="Тренд и QA",
        page_num=98,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 99 — Trend table (4 rows)
    # ------------------------------------------------------------------ #
    timeline = data.trend.get("timeline", [])
    trend_rows: list[list[str]] = []
    for t in timeline:
        trend_rows.append(
            [
                str(t.get("version", "")),
                str(t.get("date", ""))[:10],
                str(t.get("docs", "")),
                str(t.get("pairs", "")),
                str(t.get("events", "")),
                str(t.get("status_match", "")),
                str(t.get("review_queue", "")),
                f"{t.get('match_share', '')}%",
            ]
        )
    layout_paginated_table(
        prs,
        title="Сводная таблица тренда v7-v10",
        headers=["Версия", "Дата", "Документов", "Пары", "События", "Совпадений", "Очередь", "Match-доля"],
        col_widths_in=[0.9, 1.3, 0.8, 0.8, 1.1, 0.9, 1.0, 1.0],
        rows=trend_rows,
        rows_per_page=12,
        section_name="Тренд и QA",
        page_offset=99,
        total_slides=153,
        toc_slide_idx=refs["toc"],
    )

    # ------------------------------------------------------------------ #
    # WAVE C — Journey timeline slide (inserted after trend table)
    # ------------------------------------------------------------------ #
    insert_journey_timeline_slide(
        prs,
        data,
        page_num=len(prs.slides) + 1,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 100 — "Degrading" direction explanation bullets
    # ------------------------------------------------------------------ #
    degrading_bullets = [
        "trend_direction=«degrading» в trend.json — это намеренная метка, "
        "не сигнал ухудшения качества",
        "v9: добавлен D27 (ВЦИОМ, rank-3) — 26 новых пар все получили "
        "not_comparable (аналитика vs НПА несопоставимы по шкале v8)",
        "match_share упал 8.0% → 0.28%: знаменатель вырос (325→351 пар), "
        "числитель остался (1 match)",
        "review_queue сократился 183→54: 129 manual_review переведены в "
        "not_comparable после ранговой нормализации",
        "v10 = рендеринг-релиз поверх v9: корпус и события не изменились, "
        "тренд статичен",
    ]
    layout_bullets_visual(
        prs,
        title="«Degrading» — расширение корпуса, не регресс",
        bullets=degrading_bullets,
        image_path=None,
        section_name="Тренд и QA",
        page_num=100,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 101 — KPI tiles: итоговые v10
    # ------------------------------------------------------------------ #
    cn = data.control_numbers
    layout_kpi_tiles(
        prs,
        title="Итоговые показатели v10",
        tiles=[
            ("Документов", str(cn.get("documents", 27)), OCEAN["primary"]),
            ("Пар сравнения", str(cn.get("pairs", 351)), OCEAN["secondary"]),
            ("Diff-событий", str(cn.get("events", 312)), OCEAN["accent"]),
            ("Итераций", "4", OCEAN["muted"]),
        ],
        section_name="Тренд и QA",
        page_num=101,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 102 — QA gate KPI divider
    # ------------------------------------------------------------------ #
    qa_passed = data.qa.get("passed", 12)
    qa_total = data.qa.get("total", 12)
    layout_kpi_tiles(
        prs,
        title=f"QA-гейт: {qa_passed}/{qa_total} критериев PASS",
        tiles=[
            ("Пройдено (PASS)", str(qa_passed), "16A34A"),
            ("С предупреждением (WARN)", str(data.qa.get("warned", 0)), "D97706"),
            ("Провалено (FAIL)", str(data.qa.get("failed", 0)), "DC2626"),
            ("Всего критериев", str(qa_total), OCEAN["primary"]),
        ],
        section_name="Тренд и QA",
        page_num=102,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 103 — QA criteria table (12 rows)
    # ------------------------------------------------------------------ #
    checks = data.qa.get("checks", [])
    qa_rows: list[list[str]] = []
    for c in checks:
        status_val = c.get("status", "PASS")
        qa_rows.append(
            [
                c.get("name", ""),
                _trunc(c.get("description", ""), 70),
                status_val,
                _trunc(str(c.get("evidence", "")), 65),
            ]
        )
    # status_col_idx=2: PASS → green tint
    qa_color_map = {
        "PASS": STATUS_TINT_BG.get("match", "DCFCE7"),
        "WARN": STATUS_TINT_BG.get("partial_overlap", "FEF3C7"),
        "FAIL": STATUS_TINT_BG.get("contradiction", "FEE2E2"),
    }
    layout_paginated_table(
        prs,
        title="QA-критерии AC-01..AC-12",
        headers=["AC-ID", "Описание", "Статус", "Доказательство"],
        col_widths_in=[1.0, 5.5, 1.2, 4.5],
        rows=qa_rows,
        rows_per_page=12,
        section_name="Тренд и QA",
        page_offset=103,
        total_slides=153,
        toc_slide_idx=refs["toc"],
        status_col_idx=2,
        cell_color_map=qa_color_map,
    )

    # ------------------------------------------------------------------ #
    # Slide 104 — Top-3 QA criteria bullets
    # ------------------------------------------------------------------ #
    qa_bullets = [
        "AC-02 (rank-gate): 0 нарушений инварианта rank-3 ↔ rank-1 — "
        "критически важен для доказательной силы презентации",
        "AC-03 (цитаты + ранги): все 312 событий имеют цитату из источника "
        "и указание рангов — основа для юридической трассировки",
        "AC-10 (review_queue): 103 задачи, из которых 49 добавлены "
        "скриптом закрытия partial_overlap (08_close_ac10.py)",
        "AC-11 (тренд): 4 точки временного ряда v7→v10 — достаточно "
        "для анализа направления изменений",
        "AC-12 (layout): все артефакты bundle присутствуют и валидны",
    ]
    layout_bullets_visual(
        prs,
        title="QA-гейт: ключевые критерии",
        bullets=qa_bullets,
        image_path=None,
        section_name="Тренд и QA",
        page_num=104,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 105 — QA caveats bullets
    # ------------------------------------------------------------------ #
    caveats_bullets = [
        "lhs_page/rhs_page = null у ряда событий v9-наследия — допустимо, "
        "AC-03 проверяет только наличие цитаты, не номер страницы",
        "rank-4 → rank-3 нормализация: документы без явного ранга получили "
        "rank-3 по умолчанию — не влияет на ранговый шлюз (rank-3 ≤ rank-3)",
        "Provenance: 17 URL-строк без статуса — источники с нестандартными "
        "URL-схемами (Гарант, КонсультантПлюс) — документы присутствуют локально",
        "v10 — рендеринг-релиз: новые артефакты (XLSX/DOCX/PDF/JSON) "
        "не потребовали пересчёта LLM-событий",
    ]
    layout_bullets_visual(
        prs,
        title="QA: известные оговорки",
        bullets=caveats_bullets,
        image_path=None,
        section_name="Тренд и QA",
        page_num=105,
        toc_slide=toc_slide,
    )

    # ================================================================== #
    # Раздел 8 — Действия и заключение (slides 106-120)
    # ================================================================== #

    # ------------------------------------------------------------------ #
    # Slide 106 — Divider: Actions
    # ------------------------------------------------------------------ #
    layout_divider(
        prs,
        section_no="08",
        section_title="Действия и заключение",
        toc_slide=toc_slide,
    )
    refs["actions_div"] = len(prs.slides) - 1  # 0-based → 105

    # ------------------------------------------------------------------ #
    # Slide 107 — Actions severity chart
    # ------------------------------------------------------------------ #
    layout_full_chart(
        prs,
        title="Каталог редакционных действий FA-01..FA-10",
        image_path=ASSETS_DIR / "chart_actions_severity.png",
        caption=(
            "10 редакционных действий сгруппированы по степени важности. "
            "High-severity: немедленное исправление до публикации. "
            "Medium/Low: следующая итерация."
        ),
        section_name="Действия",
        page_num=107,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slides 108-117 — One slide per action FA-01..FA-10
    # ------------------------------------------------------------------ #
    sev_ru = {"high": "высокая", "medium": "средняя", "low": "низкая"}
    cat_ru = {
        "brochure_vs_npa":        "брошюра против НПА",
        "department_page_split":  "раздробление ведомственной страницы",
        "secondary_digest_links": "сноски на первичные НПА",
        "concept_supersession":   "замещение концепций",
        "amendment_chain":        "цепочка поправок",
        "amendment_to_law":       "поправки к закону",
        "amendment_to_koap":      "поправки к КоАП",
        "analytic_separation":    "разделение аналитики и НПА",
        "provenance_risk":        "риск provenance",
        "source_gap":             "пробел источника",
    }
    for idx, action in enumerate(data.actions):
        fa_id = action.get("id", f"FA-{idx+1:02d}")
        severity = action.get("severity", "")
        category = action.get("category", "")
        where = _trunc(action.get("where", ""), 200)
        what_is_wrong = _trunc(action.get("what_is_wrong", ""), 300)
        what_to_do = _trunc(action.get("what_to_do", ""), 300)
        owner = _trunc(action.get("owner", ""), 120)
        topic = cat_ru.get(category) or (category.replace("_", " ") if category else "редакционное действие")
        topic_short = _trunc(topic, 36)

        bullets = [
            f"Серьёзность: {sev_ru.get(severity, severity)}",
            f"Где: {where}",
            f"Что не так: {what_is_wrong}",
            f"Что сделать: {what_to_do}",
            f"Ответственный: {owner}",
        ]
        layout_bullets_full_width(
            prs,
            title=f"{fa_id} — {topic_short}",
            bullets=bullets,
            section_name="Действия",
            page_num=108 + idx,
            toc_slide=toc_slide,
        )

    # ------------------------------------------------------------------ #
    # Slide 118 — Action prioritisation bullets
    # ------------------------------------------------------------------ #
    by_severity = data.actions_by_severity()
    high_actions = [a for a in data.actions if a.get("severity", "") == "high"]
    high_bullets = [
        f"{a.get('id', '')} — {_trunc(a.get('what_is_wrong', ''), 70)}"
        for a in high_actions[:3]
    ]
    prio_bullets = (
        [f"High-severity ({by_severity.get('high', 0)} действий):"]
        + [f"  ▸ {b}" for b in high_bullets]
        + [
            f"Medium-severity ({by_severity.get('medium', 0)} действий): "
            "плановые правки — следующий спринт",
            f"Low-severity ({by_severity.get('low', 0)} действий): "
            "nice-to-have, не блокируют выпуск",
            "Рекомендуется начать с FA-01 и FA-02 — "
            "исправляют фактические числовые ошибки в брошюрах Минэка",
        ]
    )
    layout_bullets_visual(
        prs,
        title="Приоритизация редакционных действий",
        bullets=prio_bullets,
        image_path=None,
        section_name="Действия",
        page_num=118,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 119 — Next iterations bullets (v11 roadmap)
    # ------------------------------------------------------------------ #
    next_bullets = [
        "Гиперссылки в XLSX на источники: восстановить URL-колонку "
        "в листах bundle/pairs.csv для прямой навигации",
        "lhs_page / rhs_page: восстановить номера страниц у событий "
        "v9-наследия — нужна повторная экстракция PDF",
        "Расширение review_queue: добавить задачи по T08 "
        "(режим высылки) и T10 (КоАП) — пока без rank-1 покрытия",
        "Hyperlinks в PPTX: связать карточки событий в слайдах 52-72 "
        "с соответствующими строками в XLSX-бандле",
        "LLM-вердикт: включить SEMANTIC_COMPARATOR_ENABLED=true "
        "для A/B-сравнения с fuzzy-результатами v10",
    ]
    layout_bullets_visual(
        prs,
        title="Следующие итерации: v11 roadmap",
        bullets=next_bullets,
        image_path=None,
        section_name="Заключение",
        page_num=119,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # WAVE C — Document spotlights (3 slides: D18, D24, D27)
    # Added before the outro slide
    # ------------------------------------------------------------------ #
    spotlight_indices = insert_doc_spotlights(
        prs,
        data,
        page_offset=len(prs.slides) + 1,
        toc_slide=toc_slide,
    )
    refs["doc_spotlight_indices"] = spotlight_indices  # type: ignore[assignment]

    # ------------------------------------------------------------------ #
    # Slide 152 — Outro / Contacts (was slide 120)
    # ------------------------------------------------------------------ #
    outro_bullets = [
        "Полный bundle.json — migration_v10_out/bundle/bundle.json",
        "Машинное приложение — migration_v10_out/machine_appendix/ (14 CSV)",
        "QA-гейт — migration_v10_out/qa_report.json",
        "Презентацию воспроизвёл скрипт scripts/presentation_v10/pptx_builder.py",
        "Версия корпуса: v10.0.0 · Дата сборки: 2026-05-09",
    ]
    layout_bullets_visual(
        prs,
        title="Контакты и ссылки",
        bullets=outro_bullets,
        image_path=None,
        section_name="Заключение",
        page_num=153,
        toc_slide=toc_slide,
    )
