"""Slides 16-50: Корпус (16-25) + Pair Matrix (26-50).

Public API:
    build_part2(prs, data, *, refs) -> None

Populates refs with:
    'corpus_div'      = 15  (0-based index of slide 16)
    'pair_matrix_div' = 25  (0-based index of slide 26)
"""
from __future__ import annotations

from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

from .data_loader import V10Data
from .layout_table import layout_paginated_table
from .layouts import (
    layout_bullets_visual,
    layout_divider,
    layout_full_chart,
    layout_kpi_tiles,
)
from .pptx_helpers import add_footer
from .slides_part5 import insert_critical_pairs_section
from .theme import OCEAN, STATUS_RU, STATUS_TINT_BG

REPO_ROOT = Path(__file__).resolve().parents[3]
ASSETS_DIR = REPO_ROOT / "migration_v10_out" / "presentation" / "assets"


def _trunc(text: str, n: int) -> str:
    if len(text) <= n:
        return text
    return text[: n - 1] + "…"


def build_part2(prs: Presentation, data: V10Data, *, refs: dict[str, Any]) -> None:
    """Add slides 16-50 to *prs* and record divider indices in *refs*."""

    toc_slide = prs.slides[refs["toc"]]

    # ------------------------------------------------------------------ #
    # Slide 16 — Divider: Корпус
    # ------------------------------------------------------------------ #
    div_corpus = layout_divider(
        prs,
        section_no="02",
        section_title="Корпус: 27 документов",
        toc_slide=toc_slide,
    )
    refs["corpus_div"] = len(prs.slides) - 1  # 0-based → 15

    # ------------------------------------------------------------------ #
    # Slide 17 — KPI tiles: rank distribution
    # ------------------------------------------------------------------ #
    dbr = data.docs_by_rank()
    r1 = len(dbr.get(1, []))
    r2 = len(dbr.get(2, []))
    r3 = len(dbr.get(3, []))
    layout_kpi_tiles(
        prs,
        title="Корпус по рангам",
        tiles=[
            ("Rank-1 — официальные НПА", str(r1), OCEAN["primary"]),
            ("Rank-2 — ведомственные", str(r2), OCEAN["secondary"]),
            ("Rank-3 — аналитика", str(r3), OCEAN["accent"]),
        ],
        section_name="Корпус",
        page_num=17,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slides 18-19 — Rank-1 documents table (21 rows → 2 pages, 12+9)
    # ------------------------------------------------------------------ #
    rank1_docs = dbr.get(1, [])
    rank1_rows = [
        [
            d["id"],
            _trunc(d.get("code", ""), 18),
            _trunc(d.get("title", ""), 50),
            _trunc(d.get("type", ""), 22),
            _trunc((d.get("url") or "").split("/")[-1][:30], 30),
        ]
        for d in rank1_docs
    ]
    layout_paginated_table(
        prs,
        title="Документы rank-1 (официальные НПА)",
        headers=["ID", "Код", "Заголовок", "Тип", "Источник"],
        col_widths_in=[0.6, 1.6, 4.0, 2.0, 2.5],
        rows=rank1_rows,
        rows_per_page=12,
        section_name="Корпус",
        page_offset=18,
        total_slides=153,
        toc_slide_idx=refs["toc"],
    )

    # ------------------------------------------------------------------ #
    # Slide 20 — Rank-2 + Rank-3 combined table (6 rows, 1 page)
    # ------------------------------------------------------------------ #
    other_docs = dbr.get(2, []) + dbr.get(3, [])
    other_rows = [
        [
            d["id"],
            _trunc(d.get("code", ""), 18),
            _trunc(d.get("title", ""), 50),
            _trunc(d.get("type", ""), 22),
            _trunc((d.get("url") or "").split("/")[-1][:30], 30),
        ]
        for d in other_docs
    ]
    layout_paginated_table(
        prs,
        title="Документы rank-2 и rank-3",
        headers=["ID", "Код", "Заголовок", "Тип", "Источник"],
        col_widths_in=[0.6, 1.6, 4.0, 2.0, 2.5],
        rows=other_rows,
        rows_per_page=15,
        section_name="Корпус",
        page_offset=20,
        total_slides=153,
        toc_slide_idx=refs["toc"],
    )

    # ------------------------------------------------------------------ #
    # Slide 21 — Document types bullets
    # ------------------------------------------------------------------ #
    type_counts = Counter(d.get("type", "—") for d in data.documents)
    top5 = type_counts.most_common(5)
    type_bullets = [f"{typ}: {cnt} документов" for typ, cnt in top5]
    type_bullets.append(f"Итого: {sum(type_counts.values())} документов в корпусе")
    layout_bullets_visual(
        prs,
        title="Типы документов в корпусе",
        bullets=type_bullets,
        image_path=None,
        section_name="Корпус",
        page_num=21,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 22 — Themes distribution chart
    # ------------------------------------------------------------------ #
    layout_full_chart(
        prs,
        title="Распределение событий по темам",
        image_path=ASSETS_DIR / "chart_themes_distribution.png",
        caption=(
            "312 событий распределены по 14 темам; топ-3 темы (Инвестиционный ВНЖ, "
            "Правовое положение / патенты, ruID / цифровой въезд) покрывают ~50% массы."
        ),
        section_name="Корпус",
        page_num=22,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 23 — Theme clusters bullets (T01-T14)
    # ------------------------------------------------------------------ #
    themes_seen: dict[str, str] = {}
    for td in data.theme_doc:
        tid = td.get("theme_id", "")
        tname = td.get("theme", "")
        if tid and tid not in themes_seen:
            themes_seen[tid] = tname
    theme_bullets = [
        f"{tid}: {_trunc(tname, 60)}"
        for tid, tname in sorted(themes_seen.items())
    ]
    layout_bullets_visual(
        prs,
        title="Тематические кластеры T01-T14",
        bullets=theme_bullets,
        image_path=None,
        section_name="Корпус",
        page_num=23,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 24 — Correlation heatmap
    # ------------------------------------------------------------------ #
    layout_full_chart(
        prs,
        title="Бинарная карта: тема × документ",
        image_path=ASSETS_DIR / "chart_correlation_heatmap.png",
        caption=(
            "Бинарная карта присутствия темы в документе (14 тем × 27 документов). "
            "Белые клетки = документ не охвачен данной темой."
        ),
        section_name="Корпус",
        page_num=24,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 25 — Provenance summary bullets
    # ------------------------------------------------------------------ #
    status_counts = Counter(p.get("status", "") for p in data.provenance)
    downloaded = status_counts.get("downloaded", 0)
    copied = status_counts.get("copied", 0)
    # failures: anything not downloaded/copied/empty
    failed = sum(
        v for k, v in status_counts.items()
        if k not in ("downloaded", "copied", "")
    )
    empty = status_counts.get("", 0)
    prov_bullets = [
        f"Всего строк fetch-лога: {len(data.provenance)}",
        f"Успешно загружено (downloaded): {downloaded}",
        f"Скопировано локально (copied): {copied}",
        f"Ошибки и блокировки: {failed} (curl timeout, qrator-stub, blocked)",
        f"Без статуса / пропущено: {empty}",
    ]
    layout_bullets_visual(
        prs,
        title="Провенанс: fetch-лог источников",
        bullets=prov_bullets,
        image_path=None,
        section_name="Корпус",
        page_num=25,
        toc_slide=toc_slide,
    )

    # ================================================================== #
    # Раздел 3 — Pair Matrix (slides 26-50)
    # ================================================================== #

    # ------------------------------------------------------------------ #
    # Slide 26 — Divider: Pair Matrix
    # ------------------------------------------------------------------ #
    layout_divider(
        prs,
        section_no="03",
        section_title="Матрица пар (351)",
        toc_slide=toc_slide,
    )
    refs["pair_matrix_div"] = len(prs.slides) - 1  # 0-based → 25

    # ------------------------------------------------------------------ #
    # Slides 27-50 — 24 pages × 15 pairs (351 pairs total)
    # ------------------------------------------------------------------ #
    pairs_sorted = sorted(
        data.pairs,
        key=lambda p: int(p["id"].replace("ПАРА-", ""))
        if p["id"].replace("ПАРА-", "").isdigit()
        else 0,
    )

    pair_rows: list[list[str]] = []
    for p in pairs_sorted:
        left_short = _trunc(data.doc_short(p.get("left", "")), 22)
        right_short = _trunc(data.doc_short(p.get("right", "")), 22)
        status_eng = p.get("v8_status", "").strip()
        status_label = STATUS_RU.get(status_eng, status_eng)
        topics = _trunc(p.get("topics", ""), 30)
        pair_rows.append(
            [
                p.get("id", ""),
                left_short,
                right_short,
                topics,
                status_label,
                p.get("events_count", ""),
                p.get("rank_pair", ""),
            ]
        )

    layout_paginated_table(
        prs,
        title="Матрица пар",
        headers=["ID", "Левый", "Правый", "Темы", "Статус", "Соб.", "Ранги"],
        col_widths_in=[0.9, 1.6, 1.6, 3.5, 1.8, 0.8, 1.2],
        rows=pair_rows,
        rows_per_page=15,
        section_name="Матрица пар",
        page_offset=27,
        total_slides=153,
        toc_slide_idx=refs["toc"],
        status_col_idx=4,
        cell_color_map=STATUS_TINT_BG,
    )

    # ================================================================== #
    # WAVE C — Critical pairs section (slides 51-52)
    # Divider 03·1 + 1-page table of contradiction+partial pairs
    # ================================================================== #
    div_idx = insert_critical_pairs_section(
        prs,
        data,
        page_offset=len(prs.slides),
        toc_slide=toc_slide,
    )
    refs["critical_pairs_div"] = div_idx
