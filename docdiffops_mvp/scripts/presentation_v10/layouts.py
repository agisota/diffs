"""High-level slide layout factories for DocDiffOps v10 presentation.

Each ``layout_*`` function adds exactly one slide to *prs* and returns it.
All layouts share:
  - Light background OCEAN['bg_light']
  - Title bar (OCEAN['primary'] strip, 0.9" tall) with white 36pt bold title
  - Footer (section name | stр N/total | ↑ ToC) when page_num is provided

Layouts defined here:
  layout_cover       — full-gradient cover slide (no title bar / footer)
  layout_divider     — full-gradient section divider
  layout_kpi_tiles   — 1×4 or 2×2 grid of KPI tiles
  layout_bullets_visual — left text + optional right image
  layout_full_chart  — title + full-bleed chart image + caption
  layout_legend      — status legend grid
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from .pptx_helpers import (
    add_filled_rect,
    add_footer,
    add_gradient_bg,
    add_image,
    add_picture_centered,
    add_text_box,
    hex_to_rgb,
    set_slide_bg,
)
from .theme import (
    FONT_BODY,
    FONT_HEADER,
    FONT_MONO,
    OCEAN,
    SIZE_BULLET_PT,
    SIZE_KPI_LABEL_PT,
    SIZE_KPI_VALUE_PT,
    SIZE_TITLE_PT,
    SLIDE_H_IN,
    SLIDE_W_IN,
    TITLE_BAR_H_IN,
)

# Blank layout index (no placeholders)
_BLANK_LAYOUT_IDX = 6


def _blank_slide(prs: Presentation) -> object:
    """Add a blank slide to prs and return it."""
    layout = prs.slide_layouts[_BLANK_LAYOUT_IDX]
    return prs.slides.add_slide(layout)


def _add_title_bar(slide: object, title: str, size_pt: float | None = None) -> None:
    """Draw the primary title bar at the top of the slide.

    If size_pt is not provided, auto-selects based on title length:
    titles longer than 60 characters use 28pt to avoid overflow.
    """
    if size_pt is None:
        size_pt = 28.0 if len(title) > 60 else float(SIZE_TITLE_PT)
    add_filled_rect(slide, 0.0, 0.0, SLIDE_W_IN, TITLE_BAR_H_IN, OCEAN["primary"])
    add_text_box(
        slide,
        0.4,
        0.08,
        SLIDE_W_IN - 0.8,
        TITLE_BAR_H_IN - 0.1,
        title,
        font=FONT_HEADER,
        size_pt=size_pt,
        bold=True,
        color=OCEAN["ink_inv"],
        align="left",
        margin_zero=True,
    )


# ---------------------------------------------------------------------------
# Cover
# ---------------------------------------------------------------------------

def layout_cover(
    prs: Presentation,
    *,
    title: str,
    subtitle: str,
    date_str: str,
) -> object:
    """Full-gradient cover slide.

    Layout:
      - Gradient background (primary → secondary)
      - Large bold title, white
      - Thin accent strip above subtitle
      - Subtitle (white, slightly muted size)
      - Date string at bottom
    """
    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_dark"])

    # Full gradient rectangle
    add_gradient_bg(slide, OCEAN["primary"], OCEAN["secondary"])

    # Title
    add_text_box(
        slide,
        0.6,
        1.6,
        SLIDE_W_IN - 1.2,
        1.6,
        title,
        font=FONT_HEADER,
        size_pt=44,
        bold=True,
        color=OCEAN["ink_inv"],
        align="left",
    )

    # Thin accent strip above subtitle
    add_filled_rect(slide, 0.6, 3.35, 6.0, 0.05, OCEAN["secondary"])

    # Subtitle
    add_text_box(
        slide,
        0.6,
        3.5,
        SLIDE_W_IN - 1.2,
        1.6,
        subtitle,
        font=FONT_BODY,
        size_pt=20,
        bold=False,
        color="BFD9E8",  # white with ~70% opacity simulation via light tint
        align="left",
    )

    # Date
    add_text_box(
        slide,
        0.6,
        6.6,
        SLIDE_W_IN - 1.2,
        0.6,
        date_str,
        font=FONT_BODY,
        size_pt=14,
        bold=False,
        color="A0C4D8",
        align="left",
    )

    return slide


# ---------------------------------------------------------------------------
# Section divider
# ---------------------------------------------------------------------------

def layout_divider(
    prs: Presentation,
    *,
    section_no: str,
    section_title: str,
    prev_slide: object = None,
    next_slide: object = None,
    toc_slide: object = None,
) -> object:
    """Full-bleed gradient section divider.

    Large section number in dimmed white on the left, title center-left.
    """
    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_dark"])

    add_gradient_bg(slide, OCEAN["primary"], OCEAN["secondary"])

    # Large dimmed section number
    add_text_box(
        slide,
        0.4,
        0.8,
        4.0,
        5.5,
        section_no,
        font=FONT_HEADER,
        size_pt=180,
        bold=True,
        color="4A8EAA",  # muted blue-white at ~30% opacity simulation
        align="left",
    )

    # Section title
    add_text_box(
        slide,
        0.6,
        3.2,
        SLIDE_W_IN - 1.2,
        2.0,
        section_title,
        font=FONT_HEADER,
        size_pt=44,
        bold=True,
        color=OCEAN["ink_inv"],
        align="left",
    )

    # Navigation hints
    if prev_slide is not None:
        from .pptx_helpers import set_internal_link
        prev_tb = add_text_box(
            slide,
            0.4,
            6.8,
            3.0,
            0.5,
            "← Предыдущий раздел",
            font=FONT_BODY,
            size_pt=11,
            color="A0C4D8",
            align="left",
        )
        set_internal_link(prev_tb, prev_slide)

    if next_slide is not None:
        from .pptx_helpers import set_internal_link
        next_tb = add_text_box(
            slide,
            SLIDE_W_IN - 3.4,
            6.8,
            3.0,
            0.5,
            "Следующий раздел →",
            font=FONT_BODY,
            size_pt=11,
            color="A0C4D8",
            align="right",
        )
        set_internal_link(next_tb, next_slide)

    return slide


# ---------------------------------------------------------------------------
# KPI tiles
# ---------------------------------------------------------------------------

def layout_kpi_tiles(
    prs: Presentation,
    *,
    title: str,
    tiles: list[tuple[str, str, str]],
    section_name: str,
    page_num: int,
    toc_slide: object = None,
) -> object:
    """Slide with a grid of KPI tiles (up to 4).

    tiles: list of (label, value, color_hex) tuples.
    Layout adapts: ≤4 tiles → 1 row of 4 or 2×2 grid.
    """
    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])
    _add_title_bar(slide, title)

    n = len(tiles)
    if n <= 0:
        n = 1

    # Geometry: content area from y=1.0 to y=6.8
    content_top = 1.2
    content_h = 5.4
    margin = 0.4

    if n <= 4:
        # 2×2 grid or single row
        cols = 2 if n > 2 else n
        rows = (n + cols - 1) // cols
    else:
        cols = 4
        rows = (n + cols - 1) // cols

    tile_w = (SLIDE_W_IN - 2 * margin - (cols - 1) * 0.2) / cols
    tile_h = (content_h - (rows - 1) * 0.3) / rows

    for idx, (label, value, color_hex) in enumerate(tiles):
        col = idx % cols
        row = idx // cols
        x = margin + col * (tile_w + 0.2)
        y = content_top + row * (tile_h + 0.3)

        # Tile background
        add_filled_rect(
            slide, x, y, tile_w, tile_h, OCEAN["tile_bg"], line_hex=OCEAN["rule"]
        )

        # Value (big number)
        add_text_box(
            slide,
            x + 0.15,
            y + 0.25,
            tile_w - 0.3,
            tile_h * 0.55,
            value,
            font=FONT_HEADER,
            size_pt=SIZE_KPI_VALUE_PT,
            bold=True,
            color=color_hex if color_hex else OCEAN["primary"],
            align="center",
        )

        # Label below
        add_text_box(
            slide,
            x + 0.1,
            y + tile_h * 0.62,
            tile_w - 0.2,
            tile_h * 0.35,
            label,
            font=FONT_BODY,
            size_pt=SIZE_KPI_LABEL_PT,
            bold=False,
            color=OCEAN["ink"],
            align="center",
        )

    add_footer(
        slide,
        section_name=section_name,
        page_num=page_num,
        toc_slide=toc_slide,
    )
    return slide


# ---------------------------------------------------------------------------
# Bullets + visual
# ---------------------------------------------------------------------------

def layout_bullets_visual(
    prs: Presentation,
    *,
    title: str,
    bullets: list[str],
    image_path: Optional[Path] = None,
    section_name: str,
    page_num: int,
    toc_slide: object = None,
) -> object:
    """Two-column layout: left = bullet list, right = optional image.

    When image_path is None the text column spans full width (centered).
    """
    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])
    _add_title_bar(slide, title)

    content_top = 1.1
    content_bottom = 6.85
    content_h = content_bottom - content_top

    if image_path is not None:
        text_x = 0.5
        text_w = 5.9
        img_x = 6.7
        img_w = SLIDE_W_IN - img_x - 0.3
        img_h = content_h - 0.1
    else:
        text_x = 0.7
        text_w = SLIDE_W_IN - 1.4

    # Build bullet text
    bullet_text = "\n".join(f"• {b}" for b in bullets)
    add_text_box(
        slide,
        text_x,
        content_top + 0.15,
        text_w,
        content_h - 0.25,
        bullet_text,
        font=FONT_BODY,
        size_pt=SIZE_BULLET_PT,
        bold=False,
        color=OCEAN["ink"],
        align="left",
    )

    if image_path is not None and Path(image_path).exists():
        add_image(slide, Path(image_path), img_x, content_top, img_w, img_h)

    add_footer(
        slide,
        section_name=section_name,
        page_num=page_num,
        toc_slide=toc_slide,
    )
    return slide


# ---------------------------------------------------------------------------
# Bullets full-width (no right column)
# ---------------------------------------------------------------------------

def layout_bullets_full_width(
    prs: Presentation,
    *,
    title: str,
    bullets: list[str],
    section_name: str,
    page_num: int,
    toc_slide: object = None,
) -> object:
    """Bullets occupying the full content width (no right column).

    Used for FA-xx action slides where there is no image to show.
    Larger font (18pt) for readability across the full slide width.
    """
    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])
    _add_title_bar(slide, title)

    content_top = 1.1
    content_bottom = 6.85
    content_h = content_bottom - content_top

    bullet_text = "\n".join(f"▸ {b}" for b in bullets)
    add_text_box(
        slide,
        x_in=0.6,
        y_in=content_top + 0.15,
        w_in=SLIDE_W_IN - 1.2,
        h_in=content_h - 0.25,
        text=bullet_text,
        font=FONT_BODY,
        size_pt=18,
        bold=False,
        color=OCEAN["ink"],
        align="left",
        margin_zero=True,
    )

    add_footer(
        slide,
        section_name=section_name,
        page_num=page_num,
        toc_slide=toc_slide,
    )
    return slide


# ---------------------------------------------------------------------------
# Full chart
# ---------------------------------------------------------------------------

def layout_full_chart(
    prs: Presentation,
    *,
    title: str,
    image_path: Path,
    caption: str,
    section_name: str,
    page_num: int,
    toc_slide: object = None,
) -> object:
    """Title bar + full-width chart image + caption below."""
    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])
    _add_title_bar(slide, title)

    img_top = 1.0
    caption_h = 0.45
    footer_y = 7.0
    img_bottom = footer_y - caption_h - 0.1
    img_h = img_bottom - img_top

    img_path = Path(image_path)
    if img_path.exists():
        add_picture_centered(
            slide,
            img_path,
            content_top_in=img_top,
            max_w_in=SLIDE_W_IN - 0.6,
            max_h_in=img_h,
        )

    # Caption
    add_text_box(
        slide,
        0.5,
        footer_y - caption_h - 0.05,
        SLIDE_W_IN - 1.0,
        caption_h,
        caption,
        font=FONT_BODY,
        size_pt=10,
        bold=False,
        color=OCEAN["muted"],
        align="center",
        margin_zero=True,
    )

    add_footer(
        slide,
        section_name=section_name,
        page_num=page_num,
        toc_slide=toc_slide,
    )
    return slide


# ---------------------------------------------------------------------------
# Legend
# ---------------------------------------------------------------------------

def layout_legend(
    prs: Presentation,
    *,
    title: str,
    items: list[tuple[str, str, str]],
    section_name: str,
    page_num: int,
    toc_slide: object = None,
) -> object:
    """Legend grid: swatch square + label + glyph per row.

    items: list of (hex_color, label, glyph)
    """
    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])
    _add_title_bar(slide, title)

    content_top = 1.05
    row_h = 0.62
    swatch_size = 0.38
    row_gap = 0.1

    for i, (hex_color, label, glyph) in enumerate(items):
        y = content_top + i * (row_h + row_gap)

        # Swatch
        add_filled_rect(slide, 0.5, y + 0.1, swatch_size, swatch_size, hex_color)

        # Glyph
        add_text_box(
            slide,
            1.05,
            y + 0.05,
            0.6,
            row_h,
            glyph,
            font=FONT_HEADER,
            size_pt=22,
            bold=True,
            color=hex_color,
            align="center",
        )

        # Label
        add_text_box(
            slide,
            1.75,
            y + 0.06,
            10.0,
            row_h,
            label,
            font=FONT_BODY,
            size_pt=16,
            bold=False,
            color=OCEAN["ink"],
            align="left",
        )

    add_footer(
        slide,
        section_name=section_name,
        page_num=page_num,
        toc_slide=toc_slide,
    )
    return slide
