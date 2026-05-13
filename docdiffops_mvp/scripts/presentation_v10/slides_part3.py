"""Slides 51-95: Events (51-72) + Themes (73-80) + Review Queue (81-95).

Public API:
    build_part3(prs, data, *, refs) -> None

Populates refs with:
    'events_div' = 50   (0-based index of slide 51)
    'themes_div' = 72   (0-based index of slide 73)
    'review_div' = 80   (0-based index of slide 81)
"""
from __future__ import annotations

from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

from .data_loader import V10Data
from .layout_table import layout_paginated_table
from .layouts import (
    layout_bullets_visual,
    layout_divider,
    layout_full_chart,
    layout_kpi_tiles,
)
from .slides_part5 import (
    insert_event_detail_cards,
    insert_sankey_treemap_slides,
    insert_theme_cards,
)
from .theme import OCEAN, STATUS_RU, STATUS_TINT_BG

REPO_ROOT = Path(__file__).resolve().parents[3]
ASSETS_DIR = REPO_ROOT / "migration_v10_out" / "presentation" / "assets"


def _trunc(text: str, n: int) -> str:
    if len(text) <= n:
        return text
    return text[: n - 1] + "…"


def build_part3(prs: Presentation, data: V10Data, *, refs: dict[str, Any]) -> None:
    """Add slides 51-95 to *prs* and record divider indices in *refs*."""

    toc_slide = prs.slides[refs["toc"]]

    # ================================================================== #
    # Раздел 4 — События (slides 51-72)
    # ================================================================== #

    # ------------------------------------------------------------------ #
    # Slide 51 — Divider: Events
    # ------------------------------------------------------------------ #
    layout_divider(
        prs,
        section_no="04",
        section_title="Diff-события (312)",
        toc_slide=toc_slide,
    )
    refs["events_div"] = len(prs.slides) - 1  # 0-based → 50

    # ------------------------------------------------------------------ #
    # Slides 52-72 — 21 pages × 15 events (312 total)
    # ------------------------------------------------------------------ #
    events_sorted = sorted(
        data.events_all,
        key=lambda e: int(e["event_id"][1:])
        if e["event_id"][1:].isdigit()
        else 0,
    )

    event_rows: list[list[str]] = []
    for ev in events_sorted:
        status_eng = ev.get("status", "").strip()
        status_label = STATUS_RU.get(status_eng, status_eng)
        event_rows.append(
            [
                ev.get("event_id", ""),
                _trunc(ev.get("theme", ""), 22),
                _trunc(ev.get("left_doc", ev.get("left_id", "")), 18),
                _trunc(ev.get("right_doc", ev.get("right_id", "")), 18),
                _trunc(ev.get("claim_left", ""), 40),
                _trunc(ev.get("evidence_right", ""), 40),
                status_label,
            ]
        )

    layout_paginated_table(
        prs,
        title="Diff-события",
        headers=["ID", "Тема", "Левый", "Правый", "Утверждение", "Доказательство", "Статус"],
        col_widths_in=[0.8, 1.7, 1.0, 1.0, 3.0, 3.0, 1.8],
        rows=event_rows,
        rows_per_page=15,
        section_name="События",
        page_offset=52,
        total_slides=153,
        toc_slide_idx=refs["toc"],
        status_col_idx=6,
        cell_color_map=STATUS_TINT_BG,
        body_font_pt=8,
    )

    # ================================================================== #
    # WAVE C — Event detail cards (+11 slides: divider + 10 cards)
    # Slides 73-83 (1-based), i.e. after 21 event table pages
    # ================================================================== #
    ev_div_idx, ev_card_indices = insert_event_detail_cards(
        prs,
        data,
        page_offset=len(prs.slides),
        toc_slide=toc_slide,
    )
    refs["event_cards_div"] = ev_div_idx
    refs["event_card_indices"] = ev_card_indices  # type: ignore[assignment]

    # ================================================================== #
    # Раздел 5 — Темы и корреляция (slides 84-91 after wave C shift)
    # ================================================================== #

    # ------------------------------------------------------------------ #
    # Slide 73 — Divider: Themes
    # ------------------------------------------------------------------ #
    layout_divider(
        prs,
        section_no="05",
        section_title="Темы и корреляция",
        toc_slide=toc_slide,
    )
    refs["themes_div"] = len(prs.slides) - 1  # 0-based → 72

    # ------------------------------------------------------------------ #
    # Slide 74 — Correlation heatmap (theme × document)
    # ------------------------------------------------------------------ #
    layout_full_chart(
        prs,
        title="Корреляционная матрица: тема × документ",
        image_path=ASSETS_DIR / "chart_correlation_heatmap.png",
        caption=(
            "14 тем × 27 документов — бинарная карта присутствия. "
            "Тёмные клетки: тема охвачена документом."
        ),
        section_name="Темы и корреляция",
        page_num=74,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 75 — Coverage heatmap
    # ------------------------------------------------------------------ #
    layout_full_chart(
        prs,
        title="Глубина покрытия по рангам",
        image_path=ASSETS_DIR / "chart_coverage_heatmap.png",
        caption=(
            "Глубина покрытия по 4 рангам — где аналитика разговаривает, "
            "а официальный НПА молчит. Пустые ячейки — разрыв источника (source_gap)."
        ),
        section_name="Темы и корреляция",
        page_num=75,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 76 — Dependency graph
    # ------------------------------------------------------------------ #
    layout_full_chart(
        prs,
        title="Граф зависимостей документов",
        image_path=ASSETS_DIR / "chart_dependency_graph.png",
        caption=(
            "85 рёбер: цвета по relation_type — серые references (83), синие amends (2). "
            "Центральные узлы: документы с наибольшей степенью связности."
        ),
        section_name="Темы и корреляция",
        page_num=76,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 77 — Dependency graph observations bullets
    # ------------------------------------------------------------------ #
    rel_counts = Counter(e.get("relation_type", "") for e in data.dependency_graph)
    # Find high-degree nodes
    degree: dict[str, int] = {}
    for edge in data.dependency_graph:
        for key in ("from_doc_id", "to_doc_id"):
            node = edge.get(key, "")
            if node:
                degree[node] = degree.get(node, 0) + 1
    top_nodes = sorted(degree.items(), key=lambda x: -x[1])[:3]
    top_nodes_str = ", ".join(
        f"{data.doc_short(nid)} ({cnt} рёбер)" for nid, cnt in top_nodes
    )
    graph_bullets = [
        f"Всего рёбер в графе зависимостей: {len(data.dependency_graph)}",
        f"Тип «references»: {rel_counts.get('references', 0)} рёбер "
        f"(документы ссылаются друг на друга)",
        f"Тип «amends»: {rel_counts.get('amends', 0)} рёбра "
        f"(один акт вносит изменения в другой)",
        f"Узлы с наибольшей степенью: {_trunc(top_nodes_str, 90)}",
        "Граф симметричен — пара A↔B в матрице соответствует рёбрам A→B и B→A",
    ]
    layout_bullets_visual(
        prs,
        title="Граф зависимостей: ключевые наблюдения",
        bullets=graph_bullets,
        image_path=None,
        section_name="Темы и корреляция",
        page_num=77,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 78 — Theme catalogue table
    # ------------------------------------------------------------------ #
    themes_seen: dict[str, str] = {}
    theme_doc_count: dict[str, set] = {}
    for td in data.theme_doc:
        tid = td.get("theme_id", "")
        tname = td.get("theme", "")
        doc_id = td.get("doc_id", "")
        if tid:
            themes_seen[tid] = tname
            theme_doc_count.setdefault(tid, set()).add(doc_id)

    theme_rows = [
        [
            tid,
            _trunc(themes_seen.get(tid, ""), 45),
            str(len(theme_doc_count.get(tid, set()))),
        ]
        for tid in sorted(themes_seen.keys())
    ]
    layout_paginated_table(
        prs,
        title="Каталог тем T01-T14",
        headers=["ID", "Название темы", "Кол-во документов"],
        col_widths_in=[0.8, 9.5, 2.0],
        rows=theme_rows,
        rows_per_page=15,
        section_name="Темы и корреляция",
        page_offset=78,
        total_slides=153,
        toc_slide_idx=refs["toc"],
    )

    # ------------------------------------------------------------------ #
    # Slides 79-80 — WAVE C: Sankey + Treemap (replace weak chart+bullets)
    # ------------------------------------------------------------------ #
    insert_sankey_treemap_slides(
        prs,
        page_num_start=len(prs.slides) + 1,
        toc_slide=toc_slide,
    )

    # ================================================================== #
    # WAVE C — Theme cards (+14 slides, one per theme T01-T14)
    # Appended immediately after the themes section
    # ================================================================== #
    theme_card_indices = insert_theme_cards(
        prs,
        data,
        page_offset=len(prs.slides) + 1,
        toc_slide=toc_slide,
    )
    refs["theme_card_indices"] = theme_card_indices  # type: ignore[assignment]

    # ================================================================== #
    # Раздел 6 — Очередь ручной проверки (follows theme cards)
    # ================================================================== #

    # ------------------------------------------------------------------ #
    # Slide 81 — Divider: Review Queue
    # ------------------------------------------------------------------ #
    layout_divider(
        prs,
        section_no="06",
        section_title="Очередь ручной проверки (103 задачи)",
        toc_slide=toc_slide,
    )
    refs["review_div"] = len(prs.slides) - 1  # 0-based → 80

    # ------------------------------------------------------------------ #
    # Slide 82 — KPI tiles: priority split
    # ------------------------------------------------------------------ #
    by_priority = data.review_by_priority()
    p0 = by_priority.get("P0", 0)
    p1 = by_priority.get("P1", 0)
    p2 = by_priority.get("P2", 0)
    layout_kpi_tiles(
        prs,
        title="Очередь проверки по приоритетам",
        tiles=[
            ("P0 — критично", str(p0), "DC2626"),
            ("P1 — важно", str(p1), "D97706"),
            ("P2 — плановая", str(p2), OCEAN["secondary"]),
            ("Итого", str(p0 + p1 + p2), OCEAN["primary"]),
        ],
        section_name="Очередь проверки",
        page_num=82,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 83 — Priority split chart
    # ------------------------------------------------------------------ #
    layout_full_chart(
        prs,
        title="Приоритеты очереди ручной проверки",
        image_path=ASSETS_DIR / "chart_priority_split.png",
        caption=(
            f"P0 (критично): {p0} задачи · P1 (важно): {p1} задачи · "
            f"P2 (плановая): {p2} задач. Итого 103 задачи."
        ),
        section_name="Очередь проверки",
        page_num=83,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 84 — Review queue structure bullets
    # ------------------------------------------------------------------ #
    queue_bullets = [
        "P0 (4 задачи): критические расхождения между rank-1 актами "
        "— требуют юридического подтверждения до передачи документов",
        "P1 (2 задачи): важные несоответствия ведомственных актов "
        "официальным НПА — проверка в течение спринта",
        "P2 (97 задач): плановая сверка: частичные совпадения и устаревшие "
        "ссылки — допустимо отложить до следующей итерации",
        "Цифра 103: 54 базовых задачи (из v9) + 49 partial_overlap-событий, "
        "добавленных скриптом 08_close_ac10.py при закрытии AC-10",
        "Все задачи имеют статус «открыто» — ни одна не закрыта в рамках v10",
    ]
    layout_bullets_visual(
        prs,
        title="Структура очереди ручной проверки",
        bullets=queue_bullets,
        image_path=None,
        section_name="Очередь проверки",
        page_num=84,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 85 — TOP P0+P1 table (6 rows, 1 page)
    # ------------------------------------------------------------------ #
    p0p1_items = [r for r in data.review_queue if r.get("priority") in ("P0", "P1")]
    p0p1_rows = [
        [
            r.get("review_id", ""),
            r.get("priority", ""),
            _trunc(r.get("theme", ""), 22),
            _trunc(r.get("what_to_check", ""), 80),
            _trunc(r.get("source", ""), 22),
            _trunc(r.get("deadline", ""), 18),
            _trunc(r.get("status", ""), 12),
        ]
        for r in p0p1_items
    ]
    layout_paginated_table(
        prs,
        title="Критические задачи P0 + P1",
        headers=["RV-ID", "Приоритет", "Тема", "Что проверить", "Источник", "Дедлайн", "Статус"],
        col_widths_in=[0.9, 1.0, 1.8, 4.0, 1.7, 1.4, 1.5],
        rows=p0p1_rows,
        rows_per_page=15,
        section_name="Очередь проверки",
        page_offset=85,
        total_slides=153,
        toc_slide_idx=refs["toc"],
    )

    # ------------------------------------------------------------------ #
    # Slides 86-92 — P2 tasks (7 pages × ~14 rows)
    # ------------------------------------------------------------------ #
    p2_items = [r for r in data.review_queue if r.get("priority") == "P2"]
    p2_rows = [
        [
            r.get("review_id", ""),
            r.get("priority", ""),
            _trunc(r.get("theme", ""), 22),
            _trunc(r.get("what_to_check", ""), 80),
            _trunc(r.get("source", ""), 22),
            _trunc(r.get("deadline", ""), 18),
            _trunc(r.get("status", ""), 12),
        ]
        for r in p2_items
    ]
    layout_paginated_table(
        prs,
        title="Плановые задачи P2",
        headers=["RV-ID", "Приоритет", "Тема", "Что проверить", "Источник", "Дедлайн", "Статус"],
        col_widths_in=[0.9, 1.0, 1.8, 4.0, 1.7, 1.4, 1.5],
        rows=p2_rows,
        rows_per_page=14,
        section_name="Очередь проверки",
        page_offset=86,
        total_slides=153,
        toc_slide_idx=refs["toc"],
    )

    # ------------------------------------------------------------------ #
    # Slide 93 — Owner distribution bullets
    # ------------------------------------------------------------------ #
    owner_counts = Counter(r.get("owner", "—") for r in data.review_queue)
    top5_owners = owner_counts.most_common(5)
    owner_bullets = [
        f"{owner}: {cnt} задач{'а' if cnt == 1 else 'и' if 2 <= cnt <= 4 else ''}"
        for owner, cnt in top5_owners
    ]
    owner_bullets.append(
        f"Всего уникальных ответственных: {len(owner_counts)}"
    )
    layout_bullets_visual(
        prs,
        title="Распределение задач по ответственным",
        bullets=owner_bullets,
        image_path=None,
        section_name="Очередь проверки",
        page_num=93,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 94 — Actions severity chart
    # ------------------------------------------------------------------ #
    layout_full_chart(
        prs,
        title="Редакционные действия по степени важности",
        image_path=ASSETS_DIR / "chart_actions_severity.png",
        caption=(
            "FA-01..FA-10: 10 редакционных действий, сгруппированных "
            "по severity. High-severity требуют исправления до публикации."
        ),
        section_name="Очередь проверки",
        page_num=94,
        toc_slide=toc_slide,
    )

    # ------------------------------------------------------------------ #
    # Slide 95 — Review queue conclusions bullets
    # ------------------------------------------------------------------ #
    concl_bullets = [
        "4 задачи P0 — критические: расхождения между постановлениями "
        "Правительства требуют юридического подтверждения до передачи",
        "2 задачи P1 — ведомственные несоответствия: нарушения в блоках "
        "ЕАЭС/не-ЕАЭС и устаревшие ссылки в брошюрах",
        "97 задач P2 — плановая нагрузка: можно распределить по "
        "нескольким спринтам без блокировки публикации",
        "Основной ответственный «юрист-миграционщик» несёт наибольшую нагрузку — "
        "рекомендуется распределение на команду",
        "Ни одна задача не закрыта в v10: v10 — рендеринг-релиз, "
        "верификация запланирована на следующую итерацию",
    ]
    layout_bullets_visual(
        prs,
        title="Очередь ручной проверки: выводы",
        bullets=concl_bullets,
        image_path=None,
        section_name="Очередь проверки",
        page_num=95,
        toc_slide=toc_slide,
    )
