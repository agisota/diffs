"""Storyline for slides 1-15: Cover + ToC + Legends + Executive Summary.

Public API:
    build_part1(prs, data, *, refs) -> None

``refs`` is a dict that gets populated with section-id → 0-based slide index
so that later parts (part2, part3 …) can wire ToC hyperlinks retroactively.
Keys set here: 'cover' (0), 'toc' (2), 'executive' (5).
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from .data_loader import V10Data
from .layouts import (
    layout_bullets_visual,
    layout_cover,
    layout_divider,
    layout_full_chart,
    layout_kpi_tiles,
    layout_legend,
)
from .slides_part5 import insert_hero_stat_slide
from .pptx_helpers import (
    add_filled_rect,
    add_text_box,
    set_internal_link,
    set_slide_bg,
)
from .theme import (
    FONT_BODY,
    FONT_HEADER,
    OCEAN,
    SIZE_BULLET_PT,
    STATUS_HEX,
    STATUS_RU,
    STATUS_TO_MARK,
    V8_STATUSES,
)

REPO_ROOT = Path(__file__).resolve().parents[3]
ASSETS_DIR = REPO_ROOT / "migration_v10_out" / "presentation" / "assets"


def build_part1(prs: Presentation, data: V10Data, *, refs: dict[str, Any]) -> None:
    """Add slides 1–15 to *prs* and populate *refs* with section indices."""

    # Collect slides as we build so we can wire cross-links after creation.
    slides_added: list[object] = []

    # ------------------------------------------------------------------
    # Slide 1 — Cover
    # ------------------------------------------------------------------
    s1_cover = layout_cover(
        prs,
        title="DocDiffOps v10 — Сквозная презентация",
        subtitle=(
            "Криминалистическое сравнение корпуса нормативных и аналитических\n"
            "документов миграционной политики РФ"
        ),
        date_str="Версия 10.0.0 · 2026-05-09",
    )
    slides_added.append(s1_cover)
    refs["cover"] = 0

    # ------------------------------------------------------------------
    # Slide 2 — Abstract
    # ------------------------------------------------------------------
    s2_abstract = layout_bullets_visual(
        prs,
        title="О чём эта презентация",
        bullets=[
            "Корпус 27 документов · 351 пара · 312 diff-событий",
            "12 из 12 критериев QA-гейта пройдено (PASS)",
            "Полная сшивка 4 артефактов бандла: пояснительная записка, "
            "редакционный diff, интегральный XLSX, машинное приложение",
            "v10 — рендеринг-релиз поверх корпуса v9 (без перезапуска LLM)",
        ],
        image_path=None,
        section_name="Введение",
        page_num=2,
        toc_slide=None,  # filled in retroactively after slide 3 is created
    )
    slides_added.append(s2_abstract)

    # ------------------------------------------------------------------
    # Slide 3 — Table of Contents (ToC)
    # ------------------------------------------------------------------
    s3_toc = _build_toc_slide(prs, refs=refs)
    slides_added.append(s3_toc)
    refs["toc"] = 2

    # Back-patch slide 2 footer ToC link
    _patch_footer_toc(s2_abstract, s3_toc)

    # ------------------------------------------------------------------
    # Slide 4 — Status legend
    # ------------------------------------------------------------------
    legend_items = [
        (STATUS_HEX[s], STATUS_RU.get(s, s), STATUS_TO_MARK.get(s, ""))
        for s in V8_STATUSES
    ]
    s4_legend = layout_legend(
        prs,
        title="Легенда: 7 статусов сравнения (по шкале v8)",
        items=legend_items,
        section_name="Введение",
        page_num=4,
        toc_slide=s3_toc,
    )
    # Add footnote below the legend rows
    add_text_box(
        s4_legend,
        0.5,
        6.3,
        12.3,
        0.6,
        "Цветовая палитра наследована из docdiffops/forensic_render.py:PALETTE"
        " — соответствует XLSX/DOCX в bundle/",
        font=FONT_BODY,
        size_pt=10,
        color=OCEAN["muted"],
        align="left",
        margin_zero=True,
    )
    slides_added.append(s4_legend)

    # ------------------------------------------------------------------
    # Slide 5 — Conditions / Ranks
    # ------------------------------------------------------------------
    dbr = data.docs_by_rank()
    rank_counts = {k: len(v) for k, v in dbr.items()}
    s5_ranks = layout_bullets_visual(
        prs,
        title="Условные обозначения",
        bullets=[
            "Ранги источников:",
            "  rank-1 — официальные НПА (федеральные законы, постановления Правительства)",
            "  rank-2 — ведомственные акты и разъяснения",
            "  rank-3 — аналитика, статистика, инструкции",
            "",
            "Инвариант рангового шлюза:",
            "  rank-3 не может «опровергнуть» rank-1 — apply_rank_gate понижает"
            " такой contradiction до manual_review",
            "  Этот инвариант проверяется в QA-гейте AC-02 (PASS: 0 нарушений)",
            "",
            f"Реальное распределение: rank-1: {rank_counts.get(1, 0)} документов, "
            f"rank-2: {rank_counts.get(2, 0)}, rank-3: {rank_counts.get(3, 0)}",
        ],
        image_path=None,
        section_name="Введение",
        page_num=5,
        toc_slide=s3_toc,
    )
    slides_added.append(s5_ranks)

    # ------------------------------------------------------------------
    # Slide 6 — Section divider: Executive Summary
    # ------------------------------------------------------------------
    s6_div = layout_divider(
        prs,
        section_no="01",
        section_title="Executive Summary",
        toc_slide=s3_toc,
    )
    slides_added.append(s6_div)
    refs["executive"] = 5  # 0-based index

    # ------------------------------------------------------------------
    # Slide 7 — Hero stat (WAVE C: inserted before KPI tiles)
    # ------------------------------------------------------------------
    insert_hero_stat_slide(prs, data, page_num=7, toc_slide=s3_toc)
    slides_added.append(prs.slides[-1])

    # ------------------------------------------------------------------
    # Slide 8 — KPI tiles (was slide 7)
    # ------------------------------------------------------------------
    cn = data.control_numbers
    s8_kpi = layout_kpi_tiles(
        prs,
        title="Корпус и события",
        tiles=[
            ("Документов", str(cn.get("documents", 27)), OCEAN["primary"]),
            ("Пар сравнения", str(cn.get("pairs", 351)), OCEAN["secondary"]),
            ("Diff-событий", str(cn.get("events", 312)), OCEAN["accent"]),
            ("QA-гейт", "12/12 PASS", OCEAN["secondary"]),
        ],
        section_name="Executive Summary",
        page_num=8,
        toc_slide=s3_toc,
    )
    slides_added.append(s8_kpi)

    # ------------------------------------------------------------------
    # Slide 9 — Pair status distribution pie
    # ------------------------------------------------------------------
    s9_pie = layout_full_chart(
        prs,
        title="Распределение 351 пары по статусам",
        image_path=ASSETS_DIR / "chart_status_pie.png",
        caption=(
            "Из 351 пары: 86 на ручную проверку, 202 несопоставимы, "
            "53 частичные, 8 устаревших, 1 совпадение, 1 противоречие, 0 пробелов источника."
        ),
        section_name="Executive Summary",
        page_num=9,
        toc_slide=s3_toc,
    )
    slides_added.append(s9_pie)

    # ------------------------------------------------------------------
    # Slide 10 — Trend match_share
    # ------------------------------------------------------------------
    s10_trend = layout_full_chart(
        prs,
        title="Тренд точных совпадений по итерациям",
        image_path=ASSETS_DIR / "chart_trend_match_share.png",
        caption=(
            "Падение match_share с 8.0% до 0.28% между v8 и v9 объясняется расширением "
            "корпуса (добавлен D27 ВЦИОМ, аналитический rank-3) — "
            "это не регресс качества, а смена базы."
        ),
        section_name="Executive Summary",
        page_num=10,
        toc_slide=s3_toc,
    )
    slides_added.append(s10_trend)

    # ------------------------------------------------------------------
    # Slide 11 — Rank gate methodology
    # ------------------------------------------------------------------
    s11_rankgate = layout_bullets_visual(
        prs,
        title="Методология: ранговый шлюз (rank gate)",
        bullets=[
            "Все события классифицируются по 7-балльной шкале v8 (см. легенду на стр 4)",
            "rank-3 ↔ rank-1 события статуса contradiction → "
            "автоматическое понижение до manual_review",
            "Аналитика не может юридически опровергнуть закон — "
            "это базовое доказательственное правило",
            "AC-02 в QA-гейте проверяет: 0 нарушений шлюза в v10",
        ],
        image_path=None,
        section_name="Executive Summary",
        page_num=11,
        toc_slide=s3_toc,
    )
    slides_added.append(s11_rankgate)

    # ------------------------------------------------------------------
    # Slide 12 — Review queue trend
    # ------------------------------------------------------------------
    s12_rq = layout_full_chart(
        prs,
        title="Очередь ручной проверки по итерациям",
        image_path=ASSETS_DIR / "chart_trend_review_queue.png",
        caption=(
            "Очередь стабилизировалась на 86 событиях с v9. "
            "v10 — рендеринг-релиз, состав очереди не менялся."
        ),
        section_name="Executive Summary",
        page_num=12,
        toc_slide=s3_toc,
    )
    slides_added.append(s12_rq)

    # ------------------------------------------------------------------
    # Slide 13 — Top contradictions (was slide 12)
    # ------------------------------------------------------------------
    contrad_events = [
        e for e in data.events_all if e.get("status", "").strip() == "contradiction"
    ]
    if not contrad_events:
        contrad_events = [
            r for r in data.risks if r.get("status", "").strip() == "contradiction"
        ]

    contrad_bullets: list[str] = []
    for ev in contrad_events[:3]:
        theme = ev.get("theme", ev.get("theme_id", "—"))
        left = ev.get("left_doc", ev.get("source", "—"))
        right = ev.get("right_doc", "—")
        pair_str = f"{left} ↔ {right}"
        if len(pair_str) > 70:
            pair_str = pair_str[:67] + "…"
        claim = ev.get("claim_left", ev.get("issue", ""))[:60]
        contrad_bullets.append(f"{theme}: {pair_str}")
        if claim:
            contrad_bullets.append(f"  {claim}")

    contrad_bullets.append("Полная карточка → стр 51-72")

    s13_contrad = layout_bullets_visual(
        prs,
        title="Зафиксированные противоречия (1 на v10)",
        bullets=contrad_bullets,
        image_path=None,
        section_name="Executive Summary",
        page_num=13,
        toc_slide=s3_toc,
    )
    slides_added.append(s13_contrad)

    # ------------------------------------------------------------------
    # Slide 14 — What's new in v10 (was slide 13)
    # ------------------------------------------------------------------
    s14_new = layout_kpi_tiles(
        prs,
        title="v10 vs v9 — что изменилось",
        tiles=[
            ("Новых артефактов", "4", OCEAN["primary"]),
            ("Новых документов", "0", OCEAN["muted"]),
            ("Новых пар", "0", OCEAN["muted"]),
            ("Новых событий", "0", OCEAN["muted"]),
        ],
        section_name="Executive Summary",
        page_num=14,
        toc_slide=s3_toc,
    )
    slides_added.append(s14_new)

    # ------------------------------------------------------------------
    # Slide 15 — Top editorial actions (was slide 14)
    # ------------------------------------------------------------------
    action_bullets: list[str] = []
    for a in data.actions[:4]:
        fa_id = a.get("id", "FA-??")
        severity = a.get("severity", "")
        what = a.get("what_to_do", "")
        text = f"{fa_id} [{severity}]: {what}"
        action_bullets.append(_trunc_words(text, 80))

    s15_actions = layout_bullets_visual(
        prs,
        title="Каталог редакционных действий (FA-01..FA-10)",
        bullets=action_bullets,
        image_path=ASSETS_DIR / "chart_actions_severity.png",
        section_name="Executive Summary",
        page_num=15,
        toc_slide=s3_toc,
    )
    slides_added.append(s15_actions)

    # ------------------------------------------------------------------
    # Slide 16 — What's next (was slide 15)
    # ------------------------------------------------------------------
    s16_next = layout_bullets_visual(
        prs,
        title="Далее по презентации",
        bullets=[
            "Корпус 27 документов с разбивкой по рангам и темам — стр 17-26",
            "Полная матрица 351 пары на 24 страницах — стр 27-52",
            "Все 312 событий с цитатами + детальные карточки — стр 53-89",
            "103 задачи на ручной проверке + 10 действий — стр 99-152",
        ],
        image_path=None,
        section_name="Executive Summary",
        page_num=16,
        toc_slide=s3_toc,
    )
    slides_added.append(s16_next)

    # ------------------------------------------------------------------
    # Wire ToC hyperlinks now that all prototype slides exist
    # ------------------------------------------------------------------
    _wire_toc_links(s3_toc, slides_added, refs)


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _trunc_words(text: str, max_chars: int = 80) -> str:
    """Truncate *text* to *max_chars* on a word boundary."""
    if len(text) <= max_chars:
        return text
    cut = text[:max_chars].rsplit(" ", 1)[0]
    return cut.rstrip(",.;:") + "…"

def _build_toc_slide(prs: Presentation, *, refs: dict[str, Any]) -> object:
    """Build ToC slide (slide 3) with 8 clickable section tiles."""
    from pptx.util import Inches

    from .pptx_helpers import add_filled_rect, add_text_box
    from .layouts import _blank_slide, _add_title_bar

    slide = _blank_slide(prs)
    set_slide_bg(slide, OCEAN["bg_light"])
    _add_title_bar(slide, "Содержание")

    toc_items = [
        ("1. Executive Summary",           "стр 6-15"),
        ("2. Корпус: 27 документов",       "стр 16-25"),
        ("3. Матрица пар",                 "стр 26-50"),
        ("4. События",                     "стр 51-72"),
        ("5. Темы и корреляция",           "стр 73-80"),
        ("6. Очередь ручной проверки",     "стр 81-95"),
        ("7. Тренд и QA",                  "стр 96-105"),
        ("8. Действия и заключение",       "стр 106-120"),
        ("↑ Обложка",                      "стр 1"),
    ]

    cols = 2
    rows = (len(toc_items) + cols - 1) // cols
    margin_x = 0.4
    margin_y = 1.1
    tile_w = (13.333 - 2 * margin_x - 0.3) / cols
    tile_h = (5.6 - (rows - 1) * 0.2) / rows

    tile_shapes: list[object] = []
    for idx, (label, page_range) in enumerate(toc_items):
        col = idx % cols
        row = idx // cols
        x = margin_x + col * (tile_w + 0.3)
        y = margin_y + row * (tile_h + 0.2)

        tile = add_filled_rect(
            slide, x, y, tile_w, tile_h, OCEAN["tile_bg"], line_hex=OCEAN["rule"]
        )
        tile_shapes.append(tile)

        add_text_box(
            slide,
            x + 0.15,
            y + 0.1,
            tile_w - 0.3,
            tile_h * 0.6,
            label,
            font=FONT_HEADER,
            size_pt=16,
            bold=True,
            color=OCEAN["primary"],
            align="left",
        )
        add_text_box(
            slide,
            x + 0.15,
            y + tile_h * 0.6,
            tile_w - 0.3,
            tile_h * 0.35,
            page_range,
            font=FONT_BODY,
            size_pt=12,
            color=OCEAN["muted"],
            align="left",
        )

    # Store tile shapes in slide so we can wire links later
    slide._toc_tile_shapes = tile_shapes  # type: ignore[attr-defined]
    return slide


def _wire_toc_links(
    toc_slide: object,
    slides_added: list[object],
    refs: dict[str, Any],
) -> None:
    """Wire the 9 ToC tile shapes to their target slides.

    Tiles 0-7 (sections 1-8) → executive divider (slides_added[5]).
    Tile 8 (↑ Обложка) → cover (slides_added[0]).
    """
    tile_shapes = getattr(toc_slide, "_toc_tile_shapes", [])
    # Executive divider: cover=0, abstract=1, toc=2, legend=3, ranks=4, exec_div=5
    exec_slide = slides_added[5] if len(slides_added) > 5 else slides_added[-1]
    cover_slide = slides_added[0]

    for i, tile in enumerate(tile_shapes):
        if i == 8:
            # Last tile → cover
            set_internal_link(tile, cover_slide)
        else:
            set_internal_link(tile, exec_slide)


def _patch_footer_toc(slide: object, toc_slide: object) -> None:
    """Retroactively wire the '↑ ToC' footer shape on *slide* to *toc_slide*.

    The footer text box is the last shape added by add_footer(); we find it by
    matching its text content.
    """
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if hasattr(shape, "text_frame") and "↑ ToC" in shape.text_frame.text:
            set_internal_link(shape, toc_slide)
            break
