from __future__ import annotations

import csv
from pathlib import Path
from typing import Any

import fitz
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from .utils import compact_text, norm_text, stable_id


def make_block(doc_id: str, page_no: int | None, idx: int, text: str, block_type: str = "text", bbox: list[float] | None = None, meta: dict[str, Any] | None = None) -> dict[str, Any] | None:
    text = " ".join((text or "").split())
    if len(text) < 3:
        return None
    block_id = "blk_" + stable_id(doc_id, str(page_no), str(idx), text[:80])
    return {
        "block_id": block_id,
        "page_no": page_no,
        "block_no": idx,
        "block_type": block_type,
        "bbox": bbox,
        "text": compact_text(text, 2000),
        "norm": norm_text(text),
        "meta": meta or {},
    }


def flatten_pages(pages: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return [b for p in pages for b in p.get("blocks", [])]


def extract_pdf(path: Path, doc_id: str) -> dict[str, Any]:
    doc = fitz.open(path)
    pages = []
    for pno in range(len(doc)):
        page = doc[pno]
        blocks = []
        for i, b in enumerate(page.get_text("blocks")):
            if len(b) < 5:
                continue
            x0, y0, x1, y1, text = b[:5]
            blk = make_block(doc_id, pno + 1, i, text, bbox=[float(x0), float(y0), float(x1), float(y1)])
            if blk:
                blocks.append(blk)
        pages.append({"page_no": pno + 1, "width": page.rect.width, "height": page.rect.height, "blocks": blocks})
    return {"doc_id": doc_id, "extractor": "pymupdf_pdf", "pages": pages, "blocks": flatten_pages(pages)}


def extract_docx(path: Path, doc_id: str) -> dict[str, Any]:
    d = DocxDocument(str(path))
    blocks = []
    idx = 0
    for p in d.paragraphs:
        blk = make_block(doc_id, 1, idx, p.text, block_type="paragraph")
        if blk:
            blocks.append(blk)
            idx += 1
    for tno, table in enumerate(d.tables):
        for rno, row in enumerate(table.rows):
            text = " | ".join(cell.text for cell in row.cells)
            blk = make_block(doc_id, 1, idx, text, block_type="table_row", meta={"table": tno, "row": rno})
            if blk:
                blocks.append(blk)
                idx += 1
    pages = [{"page_no": 1, "blocks": blocks}]
    return {"doc_id": doc_id, "extractor": "python_docx", "pages": pages, "blocks": blocks}


def extract_pptx(path: Path, doc_id: str) -> dict[str, Any]:
    prs = Presentation(str(path))
    pages = []
    for sidx, slide in enumerate(prs.slides, start=1):
        blocks = []
        idx = 0
        for shape in slide.shapes:
            text = ""
            if hasattr(shape, "text"):
                text = shape.text
            if text:
                bbox = None
                try:
                    bbox = [float(shape.left), float(shape.top), float(shape.left + shape.width), float(shape.top + shape.height)]
                except Exception:
                    bbox = None
                blk = make_block(doc_id, sidx, idx, text, block_type="slide_shape", bbox=bbox)
                if blk:
                    blocks.append(blk)
                    idx += 1
        pages.append({"page_no": sidx, "blocks": blocks})
    return {"doc_id": doc_id, "extractor": "python_pptx", "pages": pages, "blocks": flatten_pages(pages)}


def extract_xlsx(path: Path, doc_id: str) -> dict[str, Any]:
    wb = load_workbook(str(path), read_only=True, data_only=False)
    pages = []
    page_no = 1
    for sheet in wb.worksheets:
        blocks = []
        idx = 0
        for row in sheet.iter_rows():
            vals = []
            coords = []
            for cell in row:
                if cell.value is not None:
                    vals.append(str(cell.value))
                    coords.append(cell.coordinate)
            if vals:
                blk = make_block(doc_id, page_no, idx, " | ".join(vals), block_type="xlsx_row", meta={"sheet": sheet.title, "cells": coords})
                if blk:
                    blocks.append(blk)
                    idx += 1
        pages.append({"page_no": page_no, "sheet": sheet.title, "blocks": blocks})
        page_no += 1
    return {"doc_id": doc_id, "extractor": "openpyxl", "pages": pages, "blocks": flatten_pages(pages)}


def extract_html(path: Path, doc_id: str) -> dict[str, Any]:
    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    texts = []
    for el in soup.find_all(["h1", "h2", "h3", "p", "li", "td", "th"]):
        tx = " ".join(el.get_text(" ").split())
        if tx:
            texts.append(tx)
    blocks = []
    for i, text in enumerate(texts):
        blk = make_block(doc_id, 1, i, text, block_type="html_block")
        if blk:
            blocks.append(blk)
    return {"doc_id": doc_id, "extractor": "beautifulsoup", "pages": [{"page_no": 1, "blocks": blocks}], "blocks": blocks}


def extract_text(path: Path, doc_id: str) -> dict[str, Any]:
    blocks = []
    text = path.read_text(encoding="utf-8", errors="ignore")
    paragraphs = [p for p in text.split("\n") if p.strip()]
    for i, p in enumerate(paragraphs):
        blk = make_block(doc_id, 1, i, p, block_type="text_line")
        if blk:
            blocks.append(blk)
    return {"doc_id": doc_id, "extractor": "plain_text", "pages": [{"page_no": 1, "blocks": blocks}], "blocks": blocks}


def extract_any(raw_path: Path, doc_id: str, canonical_pdf: Path | None = None, prefer_pdf_visual: bool = True) -> dict[str, Any]:
    # Prefer canonical PDF when available so blocks have real page bbox for red/green PDF evidence.
    if prefer_pdf_visual and canonical_pdf and canonical_pdf.exists():
        try:
            data = extract_pdf(canonical_pdf, doc_id)
            data["source_for_extraction"] = str(canonical_pdf)
            return data
        except Exception:
            pass

    ext = raw_path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf(raw_path, doc_id)
    if ext in {".docx"}:
        return extract_docx(raw_path, doc_id)
    if ext in {".pptx"}:
        return extract_pptx(raw_path, doc_id)
    if ext in {".xlsx"}:
        return extract_xlsx(raw_path, doc_id)
    if ext in {".html", ".htm"}:
        return extract_html(raw_path, doc_id)
    return extract_text(raw_path, doc_id)
